"""Inspect the gap between ThreatMetrix tables on slide 7."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = "output/DSD_Mukkamala_March_2026.pptx"
prs = Presentation(ROOT)
slide = prs.slides[6]

tables = [s for s in slide.shapes if s.has_table]
print("Tables:", len(tables))
for i, sh in enumerate(tables):
    t = sh.table
    bottom = sh.top + sh.height
    print(f"  T{i} top={sh.top} h={sh.height} bottom={bottom} rows={len(t.rows)}")
    for ri, row in enumerate(t.rows):
        print(f"    r{ri} h={row.height}")

if len(tables) >= 2:
    gap = tables[1].top - (tables[0].top + tables[0].height)
    print(f"GAP between tables: {gap} EMU ({gap/914400:.3f} in)")

for i, sh in enumerate(slide.shapes):
    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.name == "Rectangle 7":
        print(f"Rectangle7 top={sh.top} h={sh.height} bottom={sh.top+sh.height}")
        print(f"  line width={sh.line.width} fill={sh.fill.type}")

# Check graphic frame / shape outline on table shapes
for i, sh in enumerate(tables):
    el = sh.element
    ln = el.xpath(".//a:ln", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
    print(f"Table shape {i} outer ln elements:", len(ln))
    for l in ln[:3]:
        print(" ", l.attrib)
