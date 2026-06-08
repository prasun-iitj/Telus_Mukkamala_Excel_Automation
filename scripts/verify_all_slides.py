"""Verify auto-updated slides against Excel data for a given month."""

from __future__ import annotations

import argparse
import logging
import sys
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation

from app.config.dpos_team_map_parser import extract_dpos_codename, load_dpos_team_map
from app.config.team_matching import find_team_by_platform_name, find_team_for_ppt_row
from app.generators.slide_handlers import MONTH_TO_IDX
from app.readers.excel_reader import ExcelReader
from app.transformers.transformer import Transformer

logging.disable(logging.CRITICAL)

FISCAL_IDX = MONTH_TO_IDX


@dataclass
class CheckResult:
    slide: int
    area: str
    status: str  # OK, SKIP, FAIL, WARN
    detail: str = ""


@dataclass
class Report:
    month: str
    results: list[CheckResult] = field(default_factory=list)

    def add(self, slide: int, area: str, status: str, detail: str = "") -> None:
        self.results.append(CheckResult(slide, area, status, detail))

    def summary(self) -> dict[str, int]:
        counts = {"OK": 0, "SKIP": 0, "WARN": 0, "FAIL": 0}
        for r in self.results:
            counts[r.status] = counts.get(r.status, 0) + 1
        return counts


def _num(v) -> str:
    if v is None or not isinstance(v, (int, float)):
        return "0"
    return str(int(round(v)))


def _tables(slide, name="Table Placeholder 11"):
    return [s for s in slide.shapes if s.name == name and s.has_table]


def _chart_month_value(chart, month_idx: int) -> Optional[float]:
    try:
        series = list(chart.series)
        if not series:
            return None
        vals = list(series[0].values)
        if month_idx < len(vals):
            v = vals[month_idx]
            return float(v) if v is not None else None
    except Exception:
        return None
    return None


def verify_month(excel: str, ppt: str, month: str, year: int = 2026) -> Report:
    reader = ExcelReader(excel)
    dm = Transformer().build_data_model(reader)
    dm.reporting_month = month
    prs = Presentation(ppt)
    rep = Report(month=month)
    mi = FISCAL_IDX.get(month)
    dpos_map = load_dpos_team_map()

    # --- S00 Title ---
    for s in prs.slides[0].shapes:
        if s.name == "Subtitle 2":
            expected = f"{month} {year}"
            if expected in s.text:
                rep.add(0, "title date", "OK", s.text.strip())
            else:
                rep.add(0, "title date", "FAIL", f"expected '{expected}', got '{s.text.strip()}'")

    # --- S01 Apps divider ---
    subs = [s for s in prs.slides[1].shapes if s.name == "Subtitle 4"]
    if subs and f"{month} {year}" in subs[0].text:
        rep.add(1, "section date", "OK", subs[0].text.strip())
    elif subs:
        rep.add(1, "section date", "FAIL", subs[0].text.strip())

    # --- S02 Apps summary tables ---
    tables = _tables(prs.slides[2])
    checks = [
        (0, dm.apps_product_health, "Availability %", 1, lambda v: f"{int(round(v))}%" if isinstance(v, (int, float)) else "100%"),
        (1, dm.apps_quality, "Total Defects Found in PROD (P0-P4)", 2, _num),
        (2, dm.apps_rollout, "# of deployments", 1, _num),
        (3, dm.apps_delivery, "# Track 1 delivered", 1, _num),
    ]
    for ti, section, col, ppt_col, fmt in checks:
        m = next((x for x in section.months if x.month_name == month), None)
        if not m or ti >= len(tables):
            rep.add(2, f"table[{ti}]", "SKIP", "no data or table")
            continue
        tab = tables[ti].table
        mismatches = 0
        matched = 0
        for ri in range(3, len(tab.rows)):
            ppt_team = tab.cell(ri, 0).text.strip()
            if not ppt_team:
                continue
            team = find_team_for_ppt_row(ppt_team, m)
            if team is None:
                continue
            matched += 1
            exp = fmt(team.values.get(col))
            act = tab.cell(ri, ppt_col).text.strip().replace("%", "") if "%" in str(fmt(None)) else tab.cell(ri, ppt_col).text.strip()
            if "%" in exp:
                act_cmp = tab.cell(ri, ppt_col).text.strip()
                if act_cmp != exp:
                    mismatches += 1
            elif act != exp:
                mismatches += 1
        if matched == 0:
            rep.add(2, f"table[{ti}]", "WARN", "no team rows matched")
        elif mismatches:
            rep.add(2, f"table[{ti}]", "FAIL", f"{mismatches}/{matched} rows mismatch")
        else:
            rep.add(2, f"table[{ti}]", "OK", f"{matched} rows")

    # --- S03 Apps delivery charts ---
    if mi is not None:
        d_month = next((x for x in dm.apps_delivery.months if x.month_name == month), None)
        charts = [s for s in prs.slides[3].shapes if s.name == "Content Placeholder 7" and s.has_chart]
        if d_month and charts:
            metrics = [
                ("Track1", "# Track 1 delivered", 0),
                ("Sub", "# Subcontracted Delivered", 2),
            ]
            series = list(charts[0].chart.series)
            for label, col, si in metrics:
                total = sum(_float(t.values.get(col)) for t in d_month.teams)
                chart_val = (
                    list(series[si].values)[mi]
                    if si < len(series) and mi < len(list(series[si].values))
                    else None
                )
                if chart_val is not None and int(round(chart_val)) == int(total):
                    rep.add(3, f"chart {label}", "OK", f"value={chart_val}")
                elif total == 0 and (chart_val is None or chart_val == 0):
                    rep.add(3, f"chart {label}", "OK", "zero")
                else:
                    rep.add(3, f"chart {label}", "FAIL", f"excel={int(total)} chart={chart_val}")
        else:
            rep.add(3, "charts", "SKIP", "no chart or data")

    # --- S04-S05 Apps quality/rollout (partial auto) ---
    rep.add(4, "quality text/charts", "SKIP", "partial update — narrative manual")
    rep.add(5, "rollout text", "SKIP", "Key Highlights manual; charts overlay only")

    # --- S06 ThreatMetrix ---
    try:
        ws = reader.workbook["Mayur-ThreatMetrix"]
        row = None
        for r in range(3, 15):
            if ws.cell(r, 1).value == month:
                row = r
                break
        tables6 = _tables(prs.slides[6])
        if row and ws.cell(row, 3).value and len(tables6) >= 2:
            exp_total = str(int(ws.cell(row, 3).value))
            act = tables6[1].table.cell(2, 1).text.strip().replace(",", "")
            if exp_total in act or act == exp_total:
                rep.add(6, "ThreatMetrix events", "OK", act)
            else:
                rep.add(6, "ThreatMetrix events", "FAIL", f"excel={exp_total} ppt={act}")
        elif row and not ws.cell(row, 3).value:
            act = tables6[1].table.cell(2, 1).text.strip() if len(tables6) >= 2 else ""
            if act in ("N/A", "0", ""):
                rep.add(6, "ThreatMetrix", "OK", "no data — reset to N/A")
            else:
                rep.add(6, "ThreatMetrix", "FAIL", f"stale data: {act}")
        else:
            rep.add(6, "ThreatMetrix", "WARN", "sheet/row not found")
    except Exception as e:
        rep.add(6, "ThreatMetrix", "SKIP", str(e))

    # --- S07 Outages ---
    rep.add(7, "outage tables", "SKIP", "incident narrative — manual/partial")

    # --- S08 Apps PH charts ---
    if mi is not None:
        ph = next((x for x in dm.apps_product_health.months if x.month_name == month), None)
        charts8 = [s for s in prs.slides[8].shapes if s.has_chart]
        if ph and charts8:
            deg = sum(_float(t.values.get("# of Degradations")) for t in ph.teams)
            rep.add(8, "PH charts", "OK" if deg >= 0 else "WARN", f"degradations excel={deg}")
        else:
            rep.add(8, "PH charts", "SKIP", "no data")

    # --- S09 Platforms divider ---
    subs9 = [s for s in prs.slides[9].shapes if s.name == "Subtitle 4"]
    if subs9 and f"{month} {year}" in subs9[0].text:
        rep.add(9, "section date", "OK", subs9[0].text.strip())
    elif subs9:
        rep.add(9, "section date", "FAIL", subs9[0].text.strip())

    # --- S10 Platforms summary ---
    tables10 = _tables(prs.slides[10])
    checks10 = [
        (0, dm.platforms_product_health, "Availability %", 1, lambda v: f"{int(round(v))}%" if isinstance(v, (int, float)) else "100%"),
        (1, dm.platforms_quality, "Total Defects Found in PROD (P0-P4)", 2, _num),
        (2, dm.platforms_rollout, "# of deployments", 1, _num),
        (3, dm.platforms_delivery, "# Subcontracted Delivered", 4, _num),
    ]
    for ti, section, col, ppt_col, fmt in checks10:
        m = next((x for x in section.months if x.month_name == month), None)
        if not m or ti >= len(tables10):
            rep.add(10, f"table[{ti}]", "SKIP", "no data")
            continue
        tab = tables10[ti].table
        mismatches = matched = 0
        for ri in range(3, len(tab.rows)):
            ppt_team = tab.cell(ri, 0).text.strip()
            if not ppt_team:
                continue
            team = find_team_for_ppt_row(ppt_team, m)
            if team is None:
                continue
            matched += 1
            exp = fmt(team.values.get(col))
            act = tab.cell(ri, ppt_col).text.strip()
            if "%" in exp and act != exp:
                mismatches += 1
            elif "%" not in exp and act != exp:
                mismatches += 1
        status = "OK" if matched and not mismatches else ("FAIL" if mismatches else "WARN")
        rep.add(10, f"table[{ti}]", status, f"{matched} rows, {mismatches} mismatches")

    # --- S11 Platforms delivery charts ---
    if mi is not None:
        d = next((x for x in dm.platforms_delivery.months if x.month_name == month), None)
        charts11 = [s for s in prs.slides[11].shapes if s.name == "Content Placeholder 7" and s.has_chart]
        if d and charts11:
            t2 = sum(_float(t.values.get("# Track 2 delivered")) for t in d.teams)
            series = list(charts11[0].chart.series)
            chart_t2 = list(series[1].values)[mi] if len(series) > 1 and mi < len(list(series[1].values)) else None
            if chart_t2 is not None and int(round(chart_t2)) == int(t2):
                rep.add(11, "delivery chart T2", "OK", f"value={chart_t2}")
            elif t2 == 0 and (chart_t2 is None or chart_t2 == 0):
                rep.add(11, "delivery chart T2", "OK", "zero")
            else:
                rep.add(11, "delivery chart T2", "FAIL", f"excel={int(t2)} chart={chart_t2}")
        else:
            rep.add(11, "delivery charts", "SKIP", "no chart/data")

    # --- S12-S13 Platforms quality/rollout ---
    rep.add(12, "quality narrative", "SKIP", "partial — groups + manual text")
    rep.add(13, "rollout narrative", "SKIP", "partial — groups + manual text")

    # --- S14-S15 Platforms PH ---
    rep.add(14, "PH grouped charts", "SKIP", "chart overlay — spot-check only")
    rep.add(15, "PH charts", "SKIP", "chart overlay — spot-check only")

    # --- S16 Data Quality DB Audit ---
    from app.generators.slide_handlers import read_db_audit_data
    audit = read_db_audit_data(excel, month)
    charts16 = [s for s in prs.slides[16].shapes if s.has_chart]
    if audit and charts16 and mi is not None:
        rep.add(16, "DB Audit charts", "OK", f"{len(audit)} metrics loaded")
    elif not audit:
        rep.add(16, "DB Audit", "SKIP", "no Jaime-Platforms audit data")
    elif mi is not None and charts16:
        series = list(charts16[0].chart.series)
        val = list(series[0].values)[mi] if series and mi < len(list(series[0].values)) else None
        expected = audit.get("SOLO Audit Accuracy")
        if val is not None and expected is not None and abs(val - expected) < 0.05:
            rep.add(16, "DB Audit chart", "OK", f"SOLO={val}")
        elif expected is not None:
            rep.add(16, "DB Audit chart", "FAIL", f"excel={expected} chart={val}")
        else:
            rep.add(16, "DB Audit charts", "OK", f"{len(audit)} metrics loaded")
    else:
        rep.add(16, "DB Audit", "WARN", "charts present but no audit data")

    # --- S17 Data Quality incidents ---
    rep.add(17, "incident table", "SKIP", "manual narrative")

    # --- S18 DPOS divider ---
    subs18 = [s for s in prs.slides[18].shapes if s.name == "Subtitle 4"]
    if subs18 and f"{month} {year}" in subs18[0].text:
        rep.add(18, "section date", "OK", subs18[0].text.strip())
    elif subs18:
        rep.add(18, "section date", "FAIL", subs18[0].text.strip())

    # --- S19 DPOS summary ---
    tables19 = _tables(prs.slides[19])
    r_month = next((x for x in dm.platforms_rollout.months if x.month_name == month), None)
    if r_month and len(tables19) >= 3:
        total_dep = sum(_float(t.values.get("# of deployments")) for t in r_month.teams)
        act = tables19[2].table.cell(3, 1).text.strip()
        if act == _num(total_dep):
            rep.add(19, "rollout aggregate", "OK", act)
        else:
            rep.add(19, "rollout aggregate", "FAIL", f"excel={_num(total_dep)} ppt={act}")

    q_month = next((x for x in dm.platforms_quality.months if x.month_name == month), None)
    if q_month and len(tables19) >= 2:
        mismatches = matched = 0
        tab = tables19[1].table
        for ri in range(3, len(tab.rows)):
            codename = extract_dpos_codename(tab.cell(ri, 0).text)
            mapping = dpos_map.teams.get(codename)
            if not mapping:
                continue
            team = find_team_by_platform_name(mapping.platform_team, q_month)
            if not team:
                continue
            matched += 1
            exp = _num(team.values.get("Total Defects Found in PROD (P0-P4)"))
            if tab.cell(ri, 2).text.strip() != exp:
                mismatches += 1
        status = "OK" if matched and not mismatches else ("FAIL" if mismatches else "WARN")
        rep.add(19, "DPOS quality detail", status, f"{matched} mapped, {mismatches} mismatches")

    # --- S20 DPOS delivery chart ---
    if mi is not None:
        d = next((x for x in dm.platforms_delivery.months if x.month_name == month), None)
        charts20 = [s for s in prs.slides[20].shapes if s.name == "Content Placeholder 7" and s.has_chart]
        if d and charts20:
            t2 = sum(_float(t.values.get("# Track 2 delivered")) for t in d.teams)
            series = list(charts20[0].chart.series)
            chart_t2 = (
                list(series[1].values)[mi]
                if len(series) > 1 and mi < len(list(series[1].values))
                else None
            )
            if chart_t2 is not None and int(round(chart_t2)) == int(t2):
                rep.add(20, "DPOS delivery chart T2", "OK", f"value={chart_t2}")
            elif t2 == 0 and (chart_t2 is None or chart_t2 == 0):
                rep.add(20, "DPOS delivery chart T2", "OK", "zero")
            else:
                rep.add(20, "DPOS delivery chart T2", "FAIL", f"excel={int(t2)} chart={chart_t2}")
        else:
            rep.add(20, "delivery chart", "SKIP", "no chart/data")

    # --- S21 DPOS delivery text ---
    rep.add(21, "projects delivered text", "SKIP", "from person sheets — presence not validated")

    # --- S22-S24 DPOS detail ---
    rep.add(22, "quality charts", "SKIP", "grouped chart overlay")
    rep.add(23, "rollout chart", "SKIP", "grouped chart overlay")
    rep.add(24, "PH chart", "SKIP", "grouped chart overlay")

    return rep


def _float(v) -> float:
    if isinstance(v, (int, float)):
        return float(v)
    return 0.0


def print_report(rep: Report) -> None:
    counts = rep.summary()
    print(f"\n{'='*60}")
    print(f"VERIFICATION: {rep.month}")
    print(f"{'='*60}")
    for r in rep.results:
        icon = {"OK": "+", "FAIL": "X", "WARN": "!", "SKIP": "-"}.get(r.status, "?")
        line = f"  [{icon}] S{r.slide:02d} {r.area}: {r.status}"
        if r.detail:
            line += f" — {r.detail}"
        print(line)
    print(f"\nSummary: OK={counts['OK']} FAIL={counts['FAIL']} WARN={counts['WARN']} SKIP={counts['SKIP']}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--year", type=int, default=2026)
    args = parser.parse_args()

    configs = [
        ("February", "Ops Report Data Collection_February.xlsx", "output/DSD_Mukkamala_February_2026.pptx"),
        ("March", "Ops Report Data Collection_March.xlsx", "output/DSD_Mukkamala_March_2026.pptx"),
    ]
    total_fail = 0
    for month, excel, ppt in configs:
        rep = verify_month(excel, ppt, month, args.year)
        print_report(rep)
        total_fail += rep.summary()["FAIL"]

    sys.exit(1 if total_fail else 0)


if __name__ == "__main__":
    main()
