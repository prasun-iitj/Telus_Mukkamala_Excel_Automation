import zipfile, re
from openpyxl import load_workbook
from pptx import Presentation

# Excel ThreatMetrix March row
wb = load_workbook("Ops Report Data Collection_March.xlsx", data_only=True)
ws = wb["Mayur-ThreatMetrix"]
for r in range(3, 15):
    m = ws.cell(r,1).value
    if m in ("February","March"):
        vals = [ws.cell(r,c).value for c in range(1,8)]
        print(m, vals)

# PPT slide 7
prs = Presentation("output/DSD_Mukkamala_March_2026.pptx")
slide = prs.slides[6]
for sh in slide.shapes:
    if sh.has_table:
        t = sh.table
        print("TABLE", sh.name, len(t.rows), "x", len(t.columns), "h=", sh.height)
        for r in range(len(t.rows)):
            print(" ", [t.cell(r,c).text for c in range(len(t.columns))])
    elif "Arrow" in sh.name:
        print("ARROW STILL PRESENT", sh.name)

with zipfile.ZipFile("output/DSD_Mukkamala_March_2026.pptx") as z:
    xml = z.read("ppt/slides/slide7.xml").decode("utf-8")
print("XML Arrow count", xml.count("Arrow"))
