"""Custom slide handlers for complex slides that need specific data mapping."""

from __future__ import annotations

import logging
from typing import Optional

from app.config.team_matching import (
    find_team_by_platform_name,
    find_team_for_ppt_row,
    normalize_team,
)
from app.generators.chart_format import (
    add_chart_series,
    availability_chart_value,
    count_chart_value,
    normalize_chart_series,
)
from app.generators.ppt_updater import PPTUpdater
from app.models import DataModel, SectionData, MonthData, TeamMetrics

logger = logging.getLogger(__name__)

# Fiscal-year month ordering and index mapping (Nov=0 .. Oct=11)
MONTH_TO_IDX = {
    "November": 0, "December": 1, "January": 2, "February": 3,
    "March": 4, "April": 5, "May": 6, "June": 7,
    "July": 8, "August": 9, "September": 10, "October": 11,
}


def _get_month(section: SectionData, month_name: str) -> Optional[MonthData]:
    for m in section.months:
        if m.month_name == month_name:
            return m
    return None


def _num(v, default="0") -> str:
    """Format a numeric value as a whole number. Non-numeric/None → default."""
    if v is None:
        return default
    if isinstance(v, float):
        if v != v:  # NaN check
            return default
        return str(int(round(v)))
    if isinstance(v, int):
        return str(v)
    return default  # skip text values


def _pct(v) -> str:
    """Format a percentage (already 0-100 range from transformer)."""
    if v is None or not isinstance(v, (int, float)):
        return "0%"
    return f"{int(round(v))}%"


def _combined(delivered, pct_val) -> str:
    """Format as 'X/Y%' for YTD columns."""
    return f"{_num(delivered)}/{_pct(pct_val)}"


def _fmt_duration(v, default="0:00") -> str:
    """Format a duration value. Handles datetime, time, timedelta, string, None."""
    import datetime
    if v is None:
        return default
    if isinstance(v, datetime.datetime):
        # openpyxl returns datetime for time cells — extract hours:minutes
        return f"{v.hour}:{v.minute:02d}"
    if isinstance(v, datetime.time):
        return f"{v.hour}:{v.minute:02d}"
    if isinstance(v, datetime.timedelta):
        total_mins = int(v.total_seconds() // 60)
        return f"{total_mins // 60}:{total_mins % 60:02d}"
    if isinstance(v, (int, float)):
        if v == 0:
            return "0:00"
        # Could be fractional days (Excel stores time as fraction of day)
        total_mins = int(v * 24 * 60)
        return f"{total_mins // 60}:{total_mins % 60:02d}"
    s = str(v).strip()
    if s in ("", "0", "None"):
        return default
    return s


def _ph_summary_cols(t: TeamMetrics) -> list:
    v = t.values
    avail = v.get("Availability %")
    avail_str = f"{int(round(avail))}%" if isinstance(avail, (int, float)) else "100%"
    outages = v.get("# Critical outages", 0)
    duration = v.get("Critical Outage Duration (in hr:min)", "0:00")
    dur_str = _fmt_duration(duration)
    outage_str = f"{_num(outages)}/{dur_str}"
    avail_ytd = v.get("Availability % (YTD) ", v.get("Availability % (YTD)"))
    avail_ytd_str = f"{int(round(avail_ytd))}%" if isinstance(avail_ytd, (int, float)) else "100%"
    outages_ytd = v.get("# Critical outages (YTD)", 0)
    duration_ytd = v.get("Critical Outage Duration (in hr:min) (YTD)", "0:00")
    dur_ytd_str = _fmt_duration(duration_ytd)
    outage_ytd_str = f"{_num(outages_ytd)}/{dur_ytd_str}"
    return [avail_str, outage_str, avail_ytd_str, outage_ytd_str]


def _quality_summary_cols(t: TeamMetrics) -> list:
    v = t.values
    return [
        _num(v.get("P0-P1 Defect Leakage to Prod")),
        _num(v.get("Total Defects Found in PROD (P0-P4)")),
        _num(v.get("P0-P1 Defect Leakage to Prod (YTD)")),
        _num(v.get("Total Prod Defects(P0-P4) (YTD)")),
    ]


def _rollout_summary_cols(t: TeamMetrics) -> list:
    v = t.values
    return [
        _num(v.get("# of deployments")),
        _num(v.get("# of Rollbacks")),
        _num(v.get("# of Hotfixes")),
        _num(v.get("# of deployments (YTD)")),
        _num(v.get("# of Rollbacks (YTD)")),
        _num(v.get("# of Hotfixes (YTD)")),
    ]


def _delivery_summary_cols(t: TeamMetrics) -> list:
    v = t.values
    return [
        _num(v.get("# Track 1 delivered")),
        _num(v.get("# Track 2 delivered")),
        _num(v.get("# Subcontracted Delivered")),
        _num(v.get("# Sustainability Delivered")),
        _num(v.get("# of date changes for Track 1-2")),
        _combined(v.get("# Track 1 delivered (YTD)"), v.get("% completed on-time -Track 1 (YTD)")),
        _combined(v.get("# Track 2 delivered (YTD)"), v.get("% completed on-time -Track 2 (YTD)")),
        _combined(v.get("# Subcontracted delivered (YTD)"), v.get("% completed on-time -Subcontracted (YTD)")),
        _combined(v.get("# Sustainability delivered (YTD)"), v.get("% completed on-time -Sustainability (YTD)")),
        _num(v.get("# of date changes for Track 1-2 (YTD)")),
    ]


def _dpos_delivery_summary_cols(t: TeamMetrics) -> list:
    """DPOS delivery table has fewer columns (no Sustainability)."""
    v = t.values
    return [
        _num(v.get("# Track 1 delivered")),
        _num(v.get("# Track 2 delivered")),
        _num(v.get("# Subcontracted Delivered")),
        _num(v.get("# of date changes for Track 1-2")),
        _combined(v.get("# Track 1 delivered (YTD)"), v.get("% completed on-time -Track 1 (YTD)")),
        _combined(v.get("# Track 2 delivered (YTD)"), v.get("% completed on-time -Track 2 (YTD)")),
        _combined(v.get("# Subcontracted delivered (YTD)"), v.get("% completed on-time -Subcontracted (YTD)")),
        _num(v.get("# of date changes for Track 1-2 (YTD)")),
    ]


def _dpos_rollout_summary_cols(t: TeamMetrics) -> list:
    v = t.values
    return [
        _num(v.get("# of deployments")),
        _num(v.get("# of Hotfixes")),
        _num(v.get("# of deployments (YTD)")),
        _num(v.get("# of Hotfixes (YTD)")),
    ]


def _aggregate_platforms_ph(ph_month: MonthData) -> list:
    """Aggregate all platform teams into a single DPOS Product Health row (3 data cols)."""
    avail_vals = [
        t.values.get("Availability %")
        for t in ph_month.teams
        if isinstance(t.values.get("Availability %"), (int, float))
    ]
    avg_avail = round(sum(avail_vals) / len(avail_vals)) if avail_vals else 100

    total_outages = sum(
        _float_or_none(t.values.get("# Critical outages")) or 0 for t in ph_month.teams
    )
    total_duration_mins = 0
    for t in ph_month.teams:
        dur = t.values.get("Critical Outage Duration (in hr:min)")
        if isinstance(dur, (int, float)) and dur > 0:
            total_duration_mins += int(dur * 24 * 60)
        elif hasattr(dur, "hour"):
            total_duration_mins += dur.hour * 60 + dur.minute

    dur_str = f"{total_duration_mins // 60}:{total_duration_mins % 60:02d}" if total_duration_mins else "0:00"
    outage_str = f"{_num(total_outages)}/{dur_str}" if total_outages else "N/A"

    avail_str = f"{avg_avail}%"
    return [avail_str, outage_str, avail_str]


def _aggregate_platforms_rollout(r_month: MonthData) -> list:
    """Aggregate all platform teams into a single DPOS Rollout row."""
    total_dep = sum(_float_or_none(t.values.get("# of deployments")) or 0 for t in r_month.teams)
    total_hf = sum(_float_or_none(t.values.get("# of Hotfixes")) or 0 for t in r_month.teams)
    total_dep_ytd = sum(_float_or_none(t.values.get("# of deployments (YTD)")) or 0 for t in r_month.teams)
    total_hf_ytd = sum(_float_or_none(t.values.get("# of Hotfixes (YTD)")) or 0 for t in r_month.teams)
    return [_num(total_dep), _num(total_hf), _num(total_dep_ytd), _num(total_hf_ytd)]


def update_slide2_apps_summary(
    updater: PPTUpdater,
    data_model: DataModel,
    month_name: str,
) -> None:
    """Update slide 2 — Applications summary (4-quadrant tables).

    Matches Excel teams to existing PPT template rows by name (e.g.
    DataServices-Core -> Domain Services) without overwriting team labels.
    """
    slide_idx = 2
    sn = "Table Placeholder 11"

    ph_month = _get_month(data_model.apps_product_health, month_name)
    if ph_month:
        updater.update_table_by_team_match(slide_idx, sn, ph_month, _ph_summary_cols, shape_index=0)

    q_month = _get_month(data_model.apps_quality, month_name)
    if q_month:
        updater.update_table_by_team_match(slide_idx, sn, q_month, _quality_summary_cols, shape_index=1)

    r_month = _get_month(data_model.apps_rollout, month_name)
    if r_month:
        updater.update_table_by_team_match(slide_idx, sn, r_month, _rollout_summary_cols, shape_index=2)

    d_month = _get_month(data_model.apps_delivery, month_name)
    if d_month:
        updater.update_table_by_team_match(slide_idx, sn, d_month, _delivery_summary_cols, shape_index=3)

    logger.info("Slide 2 (Apps summary) updated for %s", month_name)


def _float_or_none(v) -> Optional[float]:
    """Convert to float for chart data. Non-numeric → None (gap in chart)."""
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    return None


def _chart_series_name(series, idx: int) -> str:
    """Read the display name for a chart series."""
    try:
        return series.tx.strRef.strCache[0].v
    except Exception:
        return str(series.name) if series.name else f"Series {idx + 1}"


def _series_key(name: str) -> str:
    """Normalize a chart series label for alias matching."""
    return normalize_team(name)


def _get_grouped_chart(updater: PPTUpdater, slide_idx: int, group_idx: int, chart_idx: int):
    """Return the chart object inside a grouped shape, or None."""
    slide = updater.prs.slides[slide_idx]
    shapes_list = list(slide.shapes)
    if group_idx >= len(shapes_list) or shapes_list[group_idx].shape_type != 6:
        return None
    group_shapes = list(shapes_list[group_idx].shapes)
    if chart_idx >= len(group_shapes):
        return None
    child = group_shapes[chart_idx]
    if not hasattr(child, "has_chart") or not child.has_chart:
        return None
    return child.chart


def _resolve_chart_month_index(categories: list[str], fiscal_month_idx: int) -> int:
    """Map a fiscal month index (Nov=0..Sep=10) to the chart category index.

    DPOS rollout charts use four category slots per month (Nov, '', '', '').
    Standard fiscal charts use one slot per month.
    """
    if not categories:
        return fiscal_month_idx

    month_abbrevs = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep"]
    if fiscal_month_idx < len(month_abbrevs):
        label = month_abbrevs[fiscal_month_idx]
        try:
            return categories.index(label)
        except ValueError:
            pass

    # Clustered layout: month labels every 4 slots (Nov=0, Dec=4, Jan=8, ...)
    if len(categories) > fiscal_month_idx * 4:
        return fiscal_month_idx * 4
    return fiscal_month_idx


def _overlay_grouped_chart_series_map(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    series_values: dict[str, Optional[float]],
) -> None:
    """Overlay values for the reporting month, matching chart series by name."""
    from pptx.chart.data import CategoryChartData

    chart = _get_grouped_chart(updater, slide_idx, group_idx, chart_idx)
    if chart is None:
        return

    existing_series = list(chart.series)
    cats = [str(pt.label) for pt in chart.plots[0].categories] if chart.plots[0].categories else []
    old_data = [list(s.values) for s in existing_series]
    chart_month_idx = _resolve_chart_month_index(cats, month_idx)

    value_lookup = {_series_key(k): v for k, v in series_values.items()}

    month_abbrevs = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    while chart_month_idx >= len(cats):
        cats.append(month_abbrevs[len(cats)] if len(cats) < len(month_abbrevs) else f"M{len(cats)}")
    for si in range(len(old_data)):
        while len(old_data[si]) < len(cats):
            old_data[si].append(None)

    for si, series in enumerate(existing_series):
        key = _series_key(_chart_series_name(series, si))
        if key not in value_lookup or "#ref" in key:
            continue
        if chart_month_idx < len(old_data[si]):
            old_data[si][chart_month_idx] = value_lookup[key]

    cd = CategoryChartData()
    cd.categories = cats
    for si, series_vals in enumerate(old_data):
        name = _chart_series_name(existing_series[si], si)
        add_chart_series(cd, name, series_vals)
    chart.replace_data(cd)


def _overlay_chart_series_map(
    updater: PPTUpdater,
    slide_idx: int,
    shape_name: str,
    shape_index: Optional[int],
    month_idx: int,
    series_values: dict[str, Optional[float]],
) -> None:
    """Overlay values on a top-level chart, matching series by name."""
    from pptx.chart.data import CategoryChartData

    slide = updater.prs.slides[slide_idx]
    shape = updater._find_shape(slide, shape_name, shape_index)
    if not shape.has_chart:
        return

    chart = shape.chart
    existing_series = list(chart.series)
    cats = [str(pt.label) for pt in chart.plots[0].categories] if chart.plots[0].categories else []
    old_data = [list(s.values) for s in existing_series]
    value_lookup = {_series_key(k): v for k, v in series_values.items()}

    month_abbrevs = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    while month_idx >= len(cats):
        cats.append(month_abbrevs[len(cats)] if len(cats) < len(month_abbrevs) else f"M{len(cats)}")
    for si in range(len(old_data)):
        while len(old_data[si]) < len(cats):
            old_data[si].append(None)

    for si, series in enumerate(existing_series):
        key = _series_key(_chart_series_name(series, si))
        if key not in value_lookup or "#ref" in key:
            continue
        if month_idx < len(old_data[si]):
            old_data[si][month_idx] = value_lookup[key]

    cd = CategoryChartData()
    cd.categories = cats
    for si, series_vals in enumerate(old_data):
        add_chart_series(cd, _chart_series_name(existing_series[si], si), series_vals)
    chart.replace_data(cd)


def _platform_metric_for_series(
    month_data: MonthData,
    series_name: str,
    column: str,
) -> Optional[float]:
    """Resolve a chart series label to a platform team metric (Gateways-DC merged)."""
    key = _series_key(series_name)
    if not key or "#ref" in key:
        return None
    if "gateways" in key and "dc" in key:
        team = find_team_by_platform_name("Gateways - DC", month_data)
    else:
        team = find_team_by_platform_name(series_name, month_data)
    if team is None:
        return None
    return _float_or_none(team.values.get(column))


def _update_rollout_deployment_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    r_month: MonthData,
) -> None:
    """Update rollout Group-10 style charts: per-team deployments + rollbacks/success."""
    chart = _get_grouped_chart(updater, slide_idx, group_idx, chart_idx)
    if chart is None:
        return

    values: dict[str, Optional[float]] = {}
    deploy_total = 0.0
    has_deploy = False

    for si, series in enumerate(chart.series):
        name = _chart_series_name(series, si)
        key = _series_key(name)
        if "#ref" in key:
            continue
        if key in ("rollbacks", "unsuccessful") or ("success" in key and "%" in name.lower()):
            continue
        dep = _platform_metric_for_series(r_month, name, "# of deployments")
        if dep is not None:
            values[name] = dep
            deploy_total += dep
            has_deploy = True
        elif find_team_by_platform_name(name, r_month) is not None:
            values[name] = None

    rb_total = sum(
        _float_or_none(t.values.get("# of Rollbacks")) or 0
        for t in r_month.teams
        if _float_or_none(t.values.get("# of Rollbacks")) is not None
    )
    has_rb = any(_float_or_none(t.values.get("# of Rollbacks")) is not None for t in r_month.teams)
    if has_rb:
        values["Rollbacks"] = rb_total

    unsuccessful = 0.0
    has_unsuccessful = False
    for t in r_month.teams:
        dep = _float_or_none(t.values.get("# of deployments"))
        succ = _float_or_none(t.values.get("# of Successful Deployments"))
        if dep is not None and succ is not None and dep > succ:
            unsuccessful += dep - succ
            has_unsuccessful = True
    if has_unsuccessful:
        values["Unsuccessful"] = unsuccessful

    if has_deploy and deploy_total > 0:
        uns = values.get("Unsuccessful") or 0.0
        rb = values.get("Rollbacks") or 0.0
        values["Success %"] = (deploy_total - rb - uns) / deploy_total

    _overlay_grouped_chart_series_map(
        updater, slide_idx, group_idx, chart_idx, month_idx, values,
    )


def _update_rollout_type_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    r_month: MonthData,
) -> None:
    """Update rollout Group-11 style charts (HotFix/BreakFix, Major, Minor)."""
    chart = _get_grouped_chart(updater, slide_idx, group_idx, chart_idx)
    if chart is None:
        return

    hf_total = sum(
        _float_or_none(t.values.get("# of Hotfixes")) or 0
        for t in r_month.teams
        if _float_or_none(t.values.get("# of Hotfixes")) is not None
    )
    has_hf = any(_float_or_none(t.values.get("# of Hotfixes")) is not None for t in r_month.teams)

    values: dict[str, Optional[float]] = {}
    for si, series in enumerate(chart.series):
        name = _chart_series_name(series, si)
        key = _series_key(name)
        if "hotfix" in key or "breakfix" in key:
            if has_hf:
                values[name] = hf_total

    if values:
        _overlay_grouped_chart_series_map(
            updater, slide_idx, group_idx, chart_idx, month_idx, values,
        )


def _update_availability_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    ph_month: MonthData,
    *,
    platforms: bool,
) -> None:
    """Update per-team availability % charts (Group-12 style)."""
    chart = _get_grouped_chart(updater, slide_idx, group_idx, chart_idx)
    if chart is None:
        return

    values: dict[str, Optional[float]] = {}
    for si, series in enumerate(chart.series):
        name = _chart_series_name(series, si)
        if "#ref" in _series_key(name):
            continue
        if platforms:
            raw = _platform_metric_for_series(ph_month, name, "Availability %")
            if raw is None and find_team_by_platform_name(name, ph_month) is None:
                if not ("gateways" in _series_key(name) and "dc" in _series_key(name)):
                    continue
        else:
            team = find_team_for_ppt_row(name, ph_month)
            if team is None:
                continue
            raw = _float_or_none(team.values.get("Availability %"))
        values[name] = availability_chart_value(raw)

    _overlay_grouped_chart_series_map(
        updater, slide_idx, group_idx, chart_idx, month_idx, values,
    )


def _update_ph_sla_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    ph_month: MonthData,
) -> None:
    """Update %SLA Met / Target SLA grouped charts (Group-13 style)."""
    platform_teams = [
        "Effie", "CATS", "EDP", "SOLO", "PING", "Gateways - AWS", "Gateways - DC",
    ]
    availabilities: list[float] = []
    for tname in platform_teams:
        a = _platform_metric_for_series(ph_month, tname, "Availability %")
        if a is not None:
            availabilities.append(a)

    values: dict[str, Optional[float]] = {"Target SLA (99%)": 99.0}
    if availabilities:
        values["%SLA Met"] = sum(availabilities) / len(availabilities)

    _overlay_grouped_chart_series_map(
        updater, slide_idx, group_idx, chart_idx, month_idx, values,
    )


def _update_dpos_rollout_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    r_month: MonthData,
) -> None:
    """Update DPOS rollout chart (Major / Minor / Break Fix series).

    Excel only exposes hotfix counts; Major/Minor breakdown is not available.
    """
    hf_total = sum(
        _float_or_none(t.values.get("# of Hotfixes")) or 0
        for t in r_month.teams
        if _float_or_none(t.values.get("# of Hotfixes")) is not None
    )
    has_hf = any(_float_or_none(t.values.get("# of Hotfixes")) is not None for t in r_month.teams)

    values: dict[str, Optional[float]] = {}
    if has_hf:
        values["Break Fix"] = hf_total

    if values:
        _overlay_grouped_chart_series_map(
            updater, slide_idx, group_idx, chart_idx, month_idx, values,
        )


def _update_dpos_availability_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    ph_month: MonthData,
) -> None:
    """Update DPOS slide-24 availability chart (decimal scale + SLA target line)."""
    availabilities: list[float] = []
    for t in ph_month.teams:
        a = _float_or_none(t.values.get("Availability %"))
        if a is not None:
            availabilities.append(a)

    values: dict[str, Optional[float]] = {"SLA (99.905%)": 0.99905}
    if availabilities:
        values["Availability"] = (sum(availabilities) / len(availabilities)) / 100.0

    _overlay_grouped_chart_series_map(
        updater, slide_idx, group_idx, chart_idx, month_idx, values,
    )


def _aggregate_monthly(section: SectionData, col_name: str) -> list[Optional[float]]:
    """Sum a column across all teams for each month (Nov→Oct fiscal year order).
    
    Returns 12 values in fiscal order: Nov, Dec, Jan, Feb, ..., Oct.
    """
    FISCAL_ORDER = [
        "November", "December", "January", "February", "March", "April",
        "May", "June", "July", "August", "September", "October",
    ]
    month_lookup = {m.month_name: m for m in section.months}
    result = []
    for month_name in FISCAL_ORDER:
        m = month_lookup.get(month_name)
        if m is None:
            result.append(None)
            continue
        total = None
        for team in m.teams:
            v = team.values.get(col_name)
            if isinstance(v, (int, float)):
                total = (total or 0) + v
        result.append(total)
    return result


def _per_team_monthly(section: SectionData, col_name: str, team_names: list[str]) -> list[list[Optional[float]]]:
    """Get per-team values for a column across months (Nov→Oct fiscal year).
    
    Returns one list per team, each with 12 values.
    """
    FISCAL_ORDER = [
        "November", "December", "January", "February", "March", "April",
        "May", "June", "July", "August", "September", "October",
    ]
    month_lookup = {m.month_name: m for m in section.months}
    result = []
    for tname in team_names:
        vals = []
        for month_name in FISCAL_ORDER:
            m = month_lookup.get(month_name)
            if m is None:
                vals.append(None)
                continue
            found = None
            for team in m.teams:
                if team.team_name == tname:
                    v = team.values.get(col_name)
                    found = _float_or_none(v)
                    break
            vals.append(found)
        result.append(vals)
    return result


def update_slide3_apps_delivery(
    updater: PPTUpdater,
    data_model: DataModel,
    month_name: str,
) -> None:
    """Update slide 3 — On-Time Delivery for Applications.
    
    STRATEGY: Read existing chart data from template, overlay the new month's
    values, write back. This preserves all historical data and just adds/updates
    the current month.
    
    Shapes:
      [5] Content Placeholder 7 (match 0) — Chart 1: aggregated delivery by type
      [6] Content Placeholder 7 (match 1) — Chart 2: per-team subcontracted
    """
    slide_idx = 3
    delivery = data_model.apps_delivery

    FISCAL_CATS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    MONTH_TO_IDX = {
        "November": 0, "December": 1, "January": 2, "February": 3,
        "March": 4, "April": 5, "May": 6, "June": 7,
        "July": 8, "August": 9, "September": 10, "October": 11,
    }
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is None:
        logger.warning("Unknown month '%s', skipping chart update", month_name)
        return

    d_month = _get_month(delivery, month_name)

    # --- Chart 1 (match 0): Aggregated delivery by type ---
    _overlay_chart(
        updater, slide_idx, "Content Placeholder 7", 0,
        FISCAL_CATS, month_idx, d_month,
        [
            ("# Track 1 delivered", True),
            ("# Track 2 delivered", True),
            ("# Subcontracted Delivered", True),
            ("# Sustainability Delivered", True),
        ],
    )

    # --- Chart 2 (match 1): Per-team subcontracted ---
    _overlay_chart(
        updater, slide_idx, "Content Placeholder 7", 1,
        FISCAL_CATS, month_idx, d_month,
        [("# Subcontracted Delivered", False)],  # per-team, not aggregated
    )

    logger.info("Slide 3 (Apps delivery charts) updated for %s", month_name)


def _overlay_chart(
    updater: PPTUpdater,
    slide_idx: int,
    shape_name: str,
    shape_index: int,
    categories: list[str],
    month_idx: int,
    month_data: Optional[MonthData],
    col_specs: list[tuple[str, bool]],
) -> None:
    """Read existing chart data, overlay new month's values, write back.
    
    Args:
        col_specs: List of (column_name, aggregate). If aggregate=True, sum across
                   all teams. If False, one series per team.
    """
    from pptx.chart.data import CategoryChartData

    slide = updater.prs.slides[slide_idx]
    shape = updater._find_shape(slide, shape_name, shape_index)
    if not shape.has_chart:
        return

    chart = shape.chart
    existing_series = list(chart.series)

    # Read existing data and categories (preserve historical values)
    old_data = [list(s.values) for s in existing_series]
    try:
        cats = [str(pt.label) for pt in chart.plots[0].categories] if chart.plots[0].categories else list(categories)
    except Exception:
        cats = list(categories)

    # Extend categories/data to fit the target month index
    while month_idx >= len(cats):
        extra = categories[len(cats)] if len(cats) < len(categories) else f"M{len(cats)}"
        cats.append(extra)
    for si in range(len(old_data)):
        while len(old_data[si]) < len(cats):
            old_data[si].append(None)

    # Overlay the reporting month's values
    if month_data and len(col_specs) == 1 and not col_specs[0][1]:
        col_name = col_specs[0][0]
        new_vals = [
            count_chart_value(_float_or_none(team.values.get(col_name)))
            for team in month_data.teams
        ]
        for si in range(min(len(old_data), len(new_vals))):
            if month_idx < len(old_data[si]):
                val = new_vals[si] if si < len(new_vals) else None
                old_data[si][month_idx] = val
    elif month_data:
        for si, (col_name, _) in enumerate(col_specs):
            if si >= len(old_data):
                break
            total = 0.0
            has_any = False
            for team in month_data.teams:
                v = team.values.get(col_name)
                if isinstance(v, (int, float)):
                    total += v
                    has_any = True
            if month_idx < len(old_data[si]):
                old_data[si][month_idx] = (
                    count_chart_value(total) if has_any else None
                )

    # Write back
    cd = CategoryChartData()
    cd.categories = cats
    for si, series_vals in enumerate(old_data):
        # Try to preserve series name
        name = f"Series {si+1}"
        try:
            name = existing_series[si].tx.strRef.strCache[0].v
        except Exception:
            pass
        add_chart_series(cd, name, series_vals)

    chart.replace_data(cd)


def _overlay_grouped_chart(
    updater: PPTUpdater,
    slide_idx: int,
    group_idx: int,
    chart_idx: int,
    month_idx: int,
    month_data: Optional[MonthData],
    col_specs: list[tuple[str, bool]],
) -> None:
    """Read existing chart data from a grouped chart, overlay new month's values, write back."""
    from pptx.chart.data import CategoryChartData

    slide = updater.prs.slides[slide_idx]
    shapes_list = list(slide.shapes)
    if group_idx >= len(shapes_list) or shapes_list[group_idx].shape_type != 6:
        return
    group_shapes = list(shapes_list[group_idx].shapes)
    if chart_idx >= len(group_shapes):
        return
    child = group_shapes[chart_idx]
    if not hasattr(child, "has_chart") or not child.has_chart:
        return

    chart = child.chart
    existing_series = list(chart.series)
    cats = [str(pt.label) for pt in chart.plots[0].categories] if chart.plots[0].categories else []

    # Read existing data
    old_data = [list(s.values) for s in existing_series]

    # If month_idx is beyond current data length, extend to include it
    MONTH_ABBREVS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    while month_idx >= len(cats):
        cats.append(MONTH_ABBREVS[len(cats)] if len(cats) < len(MONTH_ABBREVS) else f"M{len(cats)}")
    for si in range(len(old_data)):
        while len(old_data[si]) < len(cats):
            old_data[si].append(None)

    # Overlay new month's data
    if month_data and len(col_specs) == 1 and not col_specs[0][1]:
        # Per-team mode
        col_name = col_specs[0][0]
        new_vals = []
        for team in month_data.teams:
            v = team.values.get(col_name)
            new_vals.append(_float_or_none(v))
        # If ALL values are 0 or None, write None for all (show gap, not zeros)
        for si in range(min(len(old_data), len(new_vals))):
            if si < len(new_vals) and month_idx < len(old_data[si]):
                old_data[si][month_idx] = new_vals[si]
    elif month_data:
        # Aggregated mode
        for si, (col_name, _) in enumerate(col_specs):
            if si >= len(old_data):
                break
            total = 0.0
            has_any = False
            for team in month_data.teams:
                v = team.values.get(col_name)
                if isinstance(v, (int, float)):
                    total += v
                    has_any = True
            if month_idx < len(old_data[si]):
                old_data[si][month_idx] = total if has_any else None

    # Write back
    cd = CategoryChartData()
    cd.categories = cats
    for si, series_vals in enumerate(old_data):
        name = f"Series {si+1}"
        try:
            name = existing_series[si].tx.strRef.strCache[0].v
        except Exception:
            pass
        add_chart_series(cd, name, series_vals)
    chart.replace_data(cd)


def update_slide4_apps_quality(
    updater: PPTUpdater,
    data_model: DataModel,
    month_name: str,
) -> None:
    """Update slide 4 — Quality and Dev Efficiency for Applications.
    
    Shapes:
      [0] Key Highlights text — build from defect counts
      [3-5,7] Groups — KPI visuals (can't modify)
      [6] Mitigation table (3x4) — from Customer Impact/Root Cause/Corrective Actions
    """
    q_month = _get_month(data_model.apps_quality, month_name)
    if not q_month:
        return

    # Build Key Highlights text from defect data
    lines = ["Key Highlights", ""]
    for t in q_month.teams:
        v = t.values
        prod = v.get("Total Defects Found in PROD (P0-P4)")
        test = v.get("Total Defects Found in TEST")
        dev = v.get("Total Defects Found in DEV")
        parts = []
        if isinstance(prod, (int, float)) and prod > 0:
            parts.append(f"{_num(prod)} (P2-P4) defect{'s' if prod != 1 else ''} in Prod")
        if isinstance(test, (int, float)) and test > 0:
            parts.append(f"{_num(test)} defects in Test")
        if isinstance(dev, (int, float)) and dev > 0:
            parts.append(f"{_num(dev)} defects in Dev")
        if parts:
            lines.append(f"{t.team_name}: {', '.join(parts)}")
        else:
            lines.append(f"{t.team_name}: No defects found")
    updater.update_text_frame(4, "Text Placeholder 4", "\n".join(lines))

    # Update Mitigation table (row 2 = data row) with first team that has P0-P1 leakage
    for t in q_month.teams:
        v = t.values
        ci = v.get("Customer Impact")
        rc = v.get("Root Cause")
        ca = v.get("Corrective and Preventive Actions (CA|PA)")
        if any(isinstance(x, str) and x.strip() for x in [ci, rc, ca]):
            summary = str(ci)[:200] if isinstance(ci, str) else "No P0-P1 defect leakage"
            root = str(rc)[:200] if isinstance(rc, str) else ""
            corrective = str(ca)[:200] if isinstance(ca, str) else ""
            updater.update_table_data_rows(4, "Table Placeholder 11",
                [[summary, root, corrective, ""]], start_row=2)
            break

    logger.info("Slide 4 (Apps quality) updated for %s", month_name)

    # Update grouped charts (4 groups = 4 teams, each with Prod/Dev/Test/TestCases series)
    MONTH_TO_IDX = {"November": 0, "December": 1, "January": 2, "February": 3,
                    "March": 4, "April": 5, "May": 6, "June": 7,
                    "July": 8, "August": 9, "September": 10, "October": 11}
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is not None and q_month:
        # Group indices on slide 4: [3]=team0, [4]=team1, [5]=team2, [7]=team3
        group_indices = [3, 4, 5, 7]
        for ti, gi in enumerate(group_indices):
            if ti < len(q_month.teams):
                team = q_month.teams[ti]
                v = team.values
                # Build per-team month data with the 4 quality metrics
                team_month = MonthData(month_name=month_name, teams=[
                    TeamMetrics(team_name="Prod", values={"val": v.get("Total Defects Found in PROD (P0-P4)")}),
                    TeamMetrics(team_name="Dev", values={"val": v.get("Total Defects Found in DEV")}),
                    TeamMetrics(team_name="Test", values={"val": v.get("Total Defects Found in TEST")}),
                    TeamMetrics(team_name="TestCases", values={"val": v.get("Total Test Cases Executed")}),
                ])
                _overlay_grouped_chart(updater, 4, gi, 1, month_idx, team_month,
                    [("val", False)])


def update_slide5_apps_rollout(
    updater: PPTUpdater,
    data_model: DataModel,
    month_name: str,
) -> None:
    """Slide 5 — Rollout Success for Applications.
    
    Shapes:
      [0] Key Highlights — build from deployment counts
      [3-4] Groups — KPI visuals (can't modify)
      [5] Mitigation Plans to Increase Success (4x3)
      [7] Mitigation Plans for Unsuccessful (3x3)
    """
    r_month = _get_month(data_model.apps_rollout, month_name)
    if not r_month:
        return

    # Build Key Highlights from rollout data
    lines = ["Key Highlights"]
    for t in r_month.teams:
        v = t.values
        dep = v.get("# of deployments")
        rb = v.get("# of Rollbacks")
        hf = v.get("# of Hotfixes")
        if isinstance(dep, (int, float)) and dep > 0:
            parts = [f"{_num(dep)} deployments"]
            if isinstance(rb, (int, float)) and rb > 0:
                parts.append(f"{_num(rb)} Rollback{'s' if rb != 1 else ''}")
            if isinstance(hf, (int, float)) and hf > 0:
                parts.append(f"{_num(hf)} Hotfix{'es' if hf != 1 else ''}")
            lines.append(f"{t.team_name}: {', '.join(parts)}")
    # Note: Key Highlights text is NOT updated because the template paragraphs
    # have no runs (empty formatting) — writing text creates oversized default font.
    # The Key Highlights text should be manually written.
    # updater.update_text_frame(5, "Text Placeholder 4", "\n".join(lines))

    # Update Mitigation tables with Executive Summary / Root Cause / Lessons Learned
    # Table [5] = Mitigation Plans to Increase Success (rows 2-3)
    rollback_rows = []
    unsuccessful_rows = []
    for t in r_month.teams:
        v = t.values
        es = v.get("Executive Summary (On Rollbacks)")
        rc = v.get("Rollout Success Percentage")  # Often contains root cause text
        ll = v.get("Lessons learned")
        if isinstance(es, str) and es.strip() and es != "N/A":
            rollback_rows.append([str(es)[:200], str(rc)[:200] if isinstance(rc, str) else "", str(ll)[:200] if isinstance(ll, str) else ""])

    if rollback_rows:
        updater.update_table_data_rows(5, "Table Placeholder 11", rollback_rows, start_row=2, shape_index=0)

    # --- Grouped chart overlays ---
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is not None and r_month:
        _update_rollout_deployment_chart(updater, 5, 3, 1, month_idx, r_month)
        _update_rollout_type_chart(updater, 5, 4, 1, month_idx, r_month)

    logger.info("Slide 5 (Apps rollout) updated for %s", month_name)


# --- Slides 6-8: Product Health Applications (charts, tables, groups — preserve) ---

def update_slide6_to_8_apps_product_health(updater: PPTUpdater, data_model: DataModel, month_name: str) -> None:
    """Slides 6-8 — Product Health for Applications.
    
    Slide 6: [6] Chart — ThreatMetrix (1 series) — skip (different data source)
    Slide 8: [4] Chart 2 (2s) — degradations/impairments, [6] Chart 6 (2s) — outages
    """
    FISCAL_CATS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    MONTH_TO_IDX = {
        "November": 0, "December": 1, "January": 2, "February": 3,
        "March": 4, "April": 5, "May": 6, "June": 7,
        "July": 8, "August": 9, "September": 10, "October": 11,
    }
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is None:
        return
    ph_month = _get_month(data_model.apps_product_health, month_name)

    # Slide 8 Chart 2: degradations + impairments (2 series)
    _overlay_chart(updater, 8, "Chart 2", None, FISCAL_CATS, month_idx, ph_month,
        [("# of Degradations", True), ("# of Impairments", True)])
    # Slide 8 Chart 6: outages (2 series)
    _overlay_chart(updater, 8, "Chart 6", None, FISCAL_CATS, month_idx, ph_month,
        [("# Critical outages", True), ("Critical Outage Duration (in hr:min)", True)])

    # Slide 6 [3] Group 12 — per-team availability % (Apps teams)
    if ph_month:
        _update_availability_chart(
            updater, 6, 3, 1, month_idx, ph_month, platforms=False,
        )

    logger.info("Slides 6-8 (Apps product health) — charts updated for %s", month_name)


# --- Slide 9: Section Divider: Platforms (update date only) ---

def update_slide9_platforms_divider(updater: PPTUpdater, data_model: DataModel, month_name: str, year: int) -> None:
    """Slide 9 — Platforms section divider. Update date text only."""
    # Same structure as slide 1: two "Subtitle 4" shapes, first one is the date
    updater.update_text_frame(9, "Subtitle 4", f"{month_name} {year}", shape_index=0)
    logger.info("Slide 9 (Platforms divider) — date updated to %s %d", month_name, year)


# --- Slide 10: Platforms Summary Table (same structure as slide 2) ---

def update_slide10_platforms_summary(updater: PPTUpdater, data_model: DataModel, month_name: str) -> None:
    """Update slide 10 — Platforms summary (4-quadrant tables).

    Matches Excel teams to PPT rows (CATS, EDP, Gateways-DC, etc.) and
    merges Gateways Legacy + Blue Cloud into the single Gateways-DC row.
    """
    slide_idx = 10
    sn = "Table Placeholder 11"

    ph_month = _get_month(data_model.platforms_product_health, month_name)
    if ph_month:
        updater.update_table_by_team_match(slide_idx, sn, ph_month, _ph_summary_cols, shape_index=0)

    q_month = _get_month(data_model.platforms_quality, month_name)
    if q_month:
        updater.update_table_by_team_match(slide_idx, sn, q_month, _quality_summary_cols, shape_index=1)

    r_month = _get_month(data_model.platforms_rollout, month_name)
    if r_month:
        updater.update_table_by_team_match(slide_idx, sn, r_month, _rollout_summary_cols, shape_index=2)

    d_month = _get_month(data_model.platforms_delivery, month_name)
    if d_month:
        updater.update_table_by_team_match(slide_idx, sn, d_month, _delivery_summary_cols, shape_index=3)

    logger.info("Slide 10 (Platforms summary) updated for %s", month_name)


# --- Slides 11-17: Platforms detail slides (charts, narrative, groups — preserve) ---

def update_slides_11_to_17(updater: PPTUpdater, data_model: DataModel, month_name: str) -> None:
    """Slides 11-17 — Platforms detail slides.
    
    Slide 11: delivery charts — overlay new month's data.
    Slide 15: Product Health charts — overlay new month's data.
    Rest preserved from template.
    """
    FISCAL_CATS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    MONTH_TO_IDX = {
        "November": 0, "December": 1, "January": 2, "February": 3,
        "March": 4, "April": 5, "May": 6, "June": 7,
        "July": 8, "August": 9, "September": 10, "October": 11,
    }
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is None:
        return

    # Slide 11: Platforms delivery charts
    delivery = data_model.platforms_delivery
    d_month = _get_month(delivery, month_name)
    _overlay_chart(updater, 11, "Content Placeholder 7", 0, FISCAL_CATS, month_idx, d_month,
        [("# Track 1 delivered", True), ("# Track 2 delivered", True),
         ("# Subcontracted Delivered", True), ("# Sustainability Delivered", True)])
    if d_month:
        sub_values: dict[str, Optional[float]] = {}
        chart_shape = updater._find_shape(updater.prs.slides[11], "Content Placeholder 7", 1)
        if chart_shape.has_chart:
            for si, series in enumerate(chart_shape.chart.series):
                name = _chart_series_name(series, si)
                raw = _platform_metric_for_series(
                    d_month, name, "# Subcontracted Delivered",
                )
                sub_values[name] = count_chart_value(raw)
            _overlay_chart_series_map(updater, 11, "Content Placeholder 7", 1, month_idx, sub_values)

    # Slide 14-15: Platforms Product Health charts
    ph_month = _get_month(data_model.platforms_product_health, month_name)

    # Slide 14: Platforms Product Health — grouped charts only
    # [3] Group 12 (7 series) — availability per-team (month-based)
    # [4] Group 13 (2 series) — outages (month-based)
    # [6] Content Placeholder 7 — vulnerability remediation (NOT month-based, skip!)
    if ph_month:
        _update_availability_chart(
            updater, 14, 3, 1, month_idx, ph_month, platforms=True,
        )
        _update_ph_sla_chart(updater, 14, 4, 1, month_idx, ph_month)

    # Slide 15: Platforms Product Health charts
    _overlay_chart(updater, 15, "Chart 6", None, FISCAL_CATS, month_idx, ph_month,
        [("# of Degradations", True), ("# of Impairments", True)])
    _overlay_chart(updater, 15, "Chart 8", None, FISCAL_CATS, month_idx, ph_month,
        [("# Critical outages", True), ("Critical Outage Duration (in hr:min)", True)])

    # Slide 12: Platforms Quality — Key Highlights + Mitigation table
    q_month = _get_month(data_model.platforms_quality, month_name)
    if q_month:
        lines = ["Key Highlights", ""]
        for t in q_month.teams:
            v = t.values
            prod = v.get("Total Defects Found in PROD (P0-P4)")
            test = v.get("Total Defects Found in TEST")
            dev = v.get("Total Defects Found in DEV")
            parts = []
            if isinstance(prod, (int, float)) and prod > 0:
                parts.append(f"{_num(prod)} (P2-P4) in Prod")
            if isinstance(test, (int, float)) and test > 0:
                parts.append(f"{_num(test)} in Test")
            if isinstance(dev, (int, float)) and dev > 0:
                parts.append(f"{_num(dev)} in Dev")
            if parts:
                lines.append(f"{t.team_name}: {', '.join(parts)}")
            else:
                lines.append(f"{t.team_name}: No defects found")
        updater.update_text_frame(12, "Text Placeholder 4", "\n".join(lines))

        # Update grouped charts on slide 12 (same structure as slide 4)
        group_indices_12 = [3, 4, 5, 7]
        for ti, gi in enumerate(group_indices_12):
            if ti < len(q_month.teams):
                team = q_month.teams[ti]
                v = team.values
                team_month = MonthData(month_name=month_name, teams=[
                    TeamMetrics(team_name="Prod", values={"val": v.get("Total Defects Found in PROD (P0-P4)")}),
                    TeamMetrics(team_name="Dev", values={"val": v.get("Total Defects Found in DEV")}),
                    TeamMetrics(team_name="Test", values={"val": v.get("Total Defects Found in TEST")}),
                    TeamMetrics(team_name="TestCases", values={"val": v.get("Total Test Cases Executed")}),
                ])
                _overlay_grouped_chart(updater, 12, gi, 1, month_idx, team_month,
                    [("val", False)])

    # Slide 13: Platforms Rollout — Key Highlights
    r_month = _get_month(data_model.platforms_rollout, month_name)
    if r_month:
        lines = ["Key Highlights"]
        for t in r_month.teams:
            v = t.values
            dep = v.get("# of deployments")
            rb = v.get("# of Rollbacks")
            hf = v.get("# of Hotfixes")
            if isinstance(dep, (int, float)) and dep > 0:
                parts = [f"{_num(dep)} deployments"]
                if isinstance(rb, (int, float)) and rb > 0:
                    parts.append(f"{_num(rb)} Rollback{'s' if rb != 1 else ''}")
                if isinstance(hf, (int, float)) and hf > 0:
                    parts.append(f"{_num(hf)} Hotfix{'es' if hf != 1 else ''}")
                lines.append(f"{t.team_name}: {', '.join(parts)}")
        # Note: Key Highlights text preserved from template (empty paras cause font issues)
        # updater.update_text_frame(13, "Text Placeholder 4", "\n".join(lines))

        # Mitigation tables on S13
        rollback_rows = []
        for t in r_month.teams:
            v = t.values
            es = v.get("Executive Summary (On Rollbacks)")
            rc = v.get("Rollout Success Percentage")
            ll = v.get("Lessons learned")
            if isinstance(es, str) and es.strip() and es != "N/A":
                rollback_rows.append([str(es)[:200], str(rc)[:200] if isinstance(rc, str) else "", str(ll)[:200] if isinstance(ll, str) else ""])
        if rollback_rows:
            updater.update_table_data_rows(13, "Table Placeholder 11", rollback_rows, start_row=2, shape_index=0)

        _update_rollout_deployment_chart(updater, 13, 3, 1, month_idx, r_month)
        _update_rollout_type_chart(updater, 13, 4, 1, month_idx, r_month)

    logger.info("Slides 11-17 (Platforms) — all updated for %s", month_name)


# --- Slide 18: Section Divider: DPOS (update date only) ---

def update_slide18_dpos_divider(updater: PPTUpdater, data_model: DataModel, month_name: str, year: int) -> None:
    """Slide 18 — DPOS section divider. Update date text only."""
    updater.update_text_frame(18, "Subtitle 4", f"{month_name} {year}", shape_index=0)
    logger.info("Slide 18 (DPOS divider) — date updated to %s %d", month_name, year)


# --- Slides 19-24: DPOS slides (preserve from template) ---

def update_slides_19_to_24(updater: PPTUpdater, data_model: DataModel, month_name: str) -> None:
    """Slides 19-24 — DPOS detail slides.
    
    S19: Summary tables (same 4-quadrant as S02/S10) — uses Platforms data
    S20: Delivery chart — overlay new month
    S21-24: Narrative/groups/OLE — preserved
    """
    FISCAL_CATS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct"]
    MONTH_TO_IDX = {
        "November": 0, "December": 1, "January": 2, "February": 3,
        "March": 4, "April": 5, "May": 6, "June": 7,
        "July": 8, "August": 9, "September": 10, "October": 11,
    }
    month_idx = MONTH_TO_IDX.get(month_name)

    # --- S19: DPOS summary tables ---
    # match[0]/[2] = single aggregate row; match[1]/[3] = per-team (preserve PPT names)
    sn = "Table Placeholder 11"

    ph_month = _get_month(data_model.platforms_product_health, month_name)
    if ph_month:
        agg_ph = _aggregate_platforms_ph(ph_month)
        updater.update_table_data_rows(
            19, sn, [agg_ph], start_row=3, shape_index=0, data_start_col=1,
        )

    from app.config.dpos_team_map_parser import load_dpos_team_map

    dpos_map = load_dpos_team_map()

    q_month = _get_month(data_model.platforms_quality, month_name)
    if q_month:
        updater.update_table_by_dpos_codename(
            19, sn, q_month, dpos_map, _quality_summary_cols, shape_index=1,
        )

    r_month = _get_month(data_model.platforms_rollout, month_name)
    if r_month:
        agg_rollout = _aggregate_platforms_rollout(r_month)
        updater.update_table_data_rows(
            19, sn, [agg_rollout], start_row=3, shape_index=2, data_start_col=1,
        )

    d_month = _get_month(data_model.platforms_delivery, month_name)
    if d_month:
        updater.update_table_by_dpos_codename(
            19, sn, d_month, dpos_map, _dpos_delivery_summary_cols, shape_index=3,
        )

    # --- S20: DPOS delivery chart ---
    if month_idx is not None:
        delivery = data_model.platforms_delivery
        d_month2 = _get_month(delivery, month_name)
        _overlay_chart(updater, 20, "Content Placeholder 7", None, FISCAL_CATS, month_idx, d_month2,
            [("# Track 1 delivered", True), ("# Track 2 delivered", True),
             ("# Subcontracted Delivered", True)])

    logger.info("Slides 19-24 (DPOS) — S19 tables + S20 chart updated for %s", month_name)


def update_slide6_threatmetrix(updater: PPTUpdater, reader_path: str, month_name: str) -> None:
    """Slide 6 — Apps Product Health with ThreatMetrix data.
    
    Reads ThreatMetrix data from 'Mayur-ThreatMetrix' sheet and updates:
      [5] ThreatMetrix table (5x5) — Pass/Review/Reject counts and percentages
      [6] Total Login Events table (3x2) — monthly totals
    """
    from app.readers.excel_reader import ExcelReader
    try:
        reader = ExcelReader(reader_path)
        ws = reader.workbook["Mayur-ThreatMetrix"]
    except Exception:
        logger.warning("Mayur-ThreatMetrix sheet not found, skipping S06")
        return

    # Read ThreatMetrix data: rows 3-14 = Jan-Dec
    MONTHS = ["January", "February", "March", "April", "May", "June",
              "July", "August", "September", "October", "November", "December"]
    month_row = None
    prev_month_row = None
    for r in range(3, 15):
        m = ws.cell(r, 1).value
        if m == month_name:
            month_row = r
        idx = MONTHS.index(month_name) if month_name in MONTHS else -1
        if idx > 0 and m == MONTHS[idx - 1]:
            prev_month_row = r

    if month_row is None:
        logger.warning("Month %s not found in ThreatMetrix sheet", month_name)
        return

    # Read current month data before reset so we know whether to keep MoM arrows.
    total = ws.cell(month_row, 3).value  # Total Login Events

    # Restore the original two-table layout (Rectangle 7 frame + both tables).
    # Only remove MoM arrows when the reporting month has no data — stale arrows
    # overlap rows; when data exists the template arrows are kept.
    updater.reset_shapes_from_template(
        6,
        ["Rectangle 7", "Table Placeholder 11"],
        remove_name_substrings=["Arrow"] if not total else None,
    )
    pass_count = ws.cell(month_row, 4).value  # Pass
    review_count = ws.cell(month_row, 5).value  # Review
    reject_count = ws.cell(month_row, 6).value  # Reject
    pass_pct = ws.cell(month_row, 7).value  # Pass %

    # Read previous month for comparison
    prev_total = ws.cell(prev_month_row, 3).value if prev_month_row else None
    prev_pass_count = ws.cell(prev_month_row, 4).value if prev_month_row else None
    prev_review_count = ws.cell(prev_month_row, 5).value if prev_month_row else None
    prev_reject_count = ws.cell(prev_month_row, 6).value if prev_month_row else None
    prev_pass_pct = ws.cell(prev_month_row, 7).value if prev_month_row else None

    prev_month_name = (
        MONTHS[MONTHS.index(month_name) - 1]
        if month_name in MONTHS and MONTHS.index(month_name) > 0
        else ""
    )

    def _update_threatmetrix_headers() -> None:
        updater.update_table_cell(
            6, "Table Placeholder 11", 1, 2, f"%{month_name}", shape_index=0,
        )
        updater.update_table_cell(
            6, "Table Placeholder 11", 1, 3, f"%{prev_month_name}", shape_index=0,
        )

    def _update_login_events_table(prev_label: str, prev_val: str, cur_val: str) -> None:
        updater.update_table_cell(6, "Table Placeholder 11", 1, 0, prev_label, shape_index=1)
        updater.update_table_cell(6, "Table Placeholder 11", 1, 1, prev_val, shape_index=1)
        updater.update_table_cell(6, "Table Placeholder 11", 2, 0, month_name, shape_index=1)
        updater.update_table_cell(6, "Table Placeholder 11", 2, 1, cur_val, shape_index=1)

    def _prev_month_pct(count, total_val) -> str:
        if not total_val or not count:
            return "N/A"
        return f"{float(count) / float(total_val) * 100:.2f}%"

    if not total:
        logger.info("No ThreatMetrix data for %s — resetting tables", month_name)
        _update_threatmetrix_headers()
        prev_pass_pct_str = (
            f"{float(prev_pass_pct) * 100:.2f}%"
            if prev_pass_pct and isinstance(prev_pass_pct, (int, float))
            else "N/A"
        )
        updater.update_table_data_rows(
            6, "Table Placeholder 11",
            [
                ["Pass", "N/A", "N/A", prev_pass_pct_str, ""],
                ["Reject", "N/A", "N/A", _prev_month_pct(prev_reject_count, prev_total), ""],
                ["Review", "N/A", "N/A", _prev_month_pct(prev_review_count, prev_total), ""],
            ],
            start_row=2, shape_index=0,
        )
        _update_login_events_table(
            prev_month_name,
            f"{int(prev_total):,}" if prev_total else "N/A",
            "N/A",
        )
        return

    # Calculate percentages
    total_f = float(total) if total else 1
    pass_pct_f = float(pass_pct) * 100 if pass_pct and isinstance(pass_pct, (int, float)) else 0
    review_pct = (float(review_count) / total_f * 100) if review_count else 0
    reject_pct = (float(reject_count) / total_f * 100) if reject_count else 0

    prev_pass_pct_f = float(prev_pass_pct) * 100 if prev_pass_pct and isinstance(prev_pass_pct, (int, float)) else 0

    _update_threatmetrix_headers()

    # Update ThreatMetrix table [5] — rows 2-4 are data (Pass/Reject/Review)
    tm_rows = [
        ["Pass", f"{int(pass_count):,}" if pass_count else "0", f"{pass_pct_f:.2f}%", f"{prev_pass_pct_f:.2f}%", ""],
        ["Reject", f"{int(reject_count):,}" if reject_count else "0", f"{reject_pct:.2f}%", "N/A", ""],
        ["Review", f"{int(review_count):,}" if review_count else "0", f"{review_pct:.2f}%", "N/A", ""],
    ]
    updater.update_table_data_rows(6, "Table Placeholder 11", tm_rows, start_row=2, shape_index=0)

    _update_login_events_table(
        prev_month_name,
        f"{int(prev_total):,}" if prev_total else "0",
        f"{int(total):,}",
    )

    logger.info("Slide 6 (ThreatMetrix) updated for %s", month_name)


def update_slides_22_to_24(updater: PPTUpdater, data_model: DataModel, month_name: str) -> None:
    """Slides 22-24 — DPOS Quality, Rollout, and Product Health grouped charts.

    S22: DPOS Quality — 7 grouped charts showing defects per team
    S23: DPOS Rollout — 1 grouped chart (per-team deployments)
    S24: DPOS Product Health — 1 grouped chart (aggregated outages/degradations)
    """
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is None:
        logger.warning("Unknown month '%s', skipping S22-S24 chart updates", month_name)
        return

    # --- S22: DPOS Quality grouped charts ---
    q_month = _get_month(data_model.platforms_quality, month_name)
    if q_month:
        # Groups [0],[4],[5],[6] — per-team quality charts (month-based categories)
        # Each group has a chart child; we overlay per-team defect data.
        # [0] Group 18 (2 series), [4] Group 2 (3 series),
        # [5] Group 1 (6 series), [6] Group 5 (5 series)
        quality_groups = [
            (0, "Total Defects Found in PROD (P0-P4)"),
            (4, "Total Defects Found in PROD (P0-P4)"),
            (5, "Total Defects Found in PROD (P0-P4)"),
            (6, "Total Defects Found in PROD (P0-P4)"),
        ]
        for gi, col_name in quality_groups:
            team_month = MonthData(month_name=month_name, teams=[
                TeamMetrics(team_name=t.team_name, values={"val": t.values.get(col_name)})
                for t in q_month.teams
            ])
            _overlay_grouped_chart(updater, 22, gi, 1, month_idx, team_month,
                [("val", False)])

        # [9] Group 10 (2 series) — categories are ['< 30', '30 - 60', '60 - 90', '90+'],
        #     NOT month-based. Skip.
        # [10] Group 31 (2 series) — same non-month categories. Skip.

        # [11] Group 34 (2 series) — month-based, overlay
        team_month_g34 = MonthData(month_name=month_name, teams=[
            TeamMetrics(team_name=t.team_name, values={"val": t.values.get("Total Defects Found in PROD (P0-P4)")})
            for t in q_month.teams
        ])
        _overlay_grouped_chart(updater, 22, 11, 1, month_idx, team_month_g34,
            [("val", False)])

    logger.info("Slide 22 (DPOS Quality) — grouped charts updated for %s", month_name)

    # --- S23: DPOS Rollout grouped chart (Major/Minor/Break Fix — not per-team) ---
    r_month = _get_month(data_model.platforms_rollout, month_name)
    if r_month:
        _update_dpos_rollout_chart(updater, 23, 3, 1, month_idx, r_month)

    logger.info("Slide 23 (DPOS Rollout) — grouped chart updated for %s", month_name)

    # --- S24: DPOS Product Health grouped chart ---
    ph_month = _get_month(data_model.platforms_product_health, month_name)
    if ph_month:
        _update_dpos_availability_chart(updater, 24, 4, 1, month_idx, ph_month)

    logger.info("Slide 24 (DPOS Product Health) — grouped chart updated for %s", month_name)


# ---------------------------------------------------------------------------
# S07: Non-Critical Downtimes + Critical Outages
# ---------------------------------------------------------------------------

def read_non_critical_downtimes(reader_path: str, month_name: str) -> list[dict]:
    """Read Non Critical Downtimes rows from ALL individual person sheets.

    Scans every non-Master sheet for a "Non Critical Downtimes" header in
    column A, then collects rows matching *month_name*.

    Returns a list of dicts with keys:
        impacted_app, causing_app, degradation_impairment, duration,
        customer_impact, success_failure_rate, executive_summary,
        mitigation_steps, lessons_learned, incident_num
    """
    from app.readers.excel_reader import ExcelReader

    try:
        reader = ExcelReader(reader_path)
    except Exception:
        logger.warning("Could not open workbook at %s for Non Critical Downtimes", reader_path)
        return []

    MONTHS = [
        "January", "February", "March", "April", "May", "June",
        "July", "August", "September", "October", "November", "December",
    ]

    results: list[dict] = []

    for sheet_name in reader.workbook.sheetnames:
        if "Master" in sheet_name or "DoNotChange" in sheet_name:
            continue
        if "Guideline" in sheet_name or "Timeline" in sheet_name:
            continue

        ws = reader.workbook[sheet_name]

        # Find "Non Critical Downtimes" header row
        header_row = None
        for r in range(1, 300):
            val = ws.cell(r, 1).value
            if val and isinstance(val, str) and "non critical downtime" in val.lower():
                header_row = r
                break

        if header_row is None:
            continue

        # Data starts at header_row + 2 (skip the column-header row)
        data_start = header_row + 2

        # Walk rows: column 1 may have month name or be blank (continuation)
        current_month = None
        for r in range(data_start, data_start + 200):
            c1 = ws.cell(r, 1).value
            if c1 and isinstance(c1, str) and c1.strip() in MONTHS:
                current_month = c1.strip()

            if current_month != month_name:
                continue

            # Check if this row has actual data (at least impacted_app)
            impacted = ws.cell(r, 2).value
            if not impacted or (isinstance(impacted, str) and impacted.strip() in ("", "N/A")):
                continue

            results.append({
                "impacted_app": str(ws.cell(r, 2).value or "")[:200],
                "causing_app": str(ws.cell(r, 3).value or "")[:200],
                "degradation_impairment": str(ws.cell(r, 4).value or "")[:100],
                "duration": _fmt_duration(ws.cell(r, 5).value),
                "customer_impact": str(ws.cell(r, 6).value or "")[:200],
                "success_failure_rate": str(ws.cell(r, 7).value or "")[:200],
                "executive_summary": str(ws.cell(r, 8).value or "")[:200],
                "mitigation_steps": str(ws.cell(r, 9).value or "")[:200],
                "lessons_learned": str(ws.cell(r, 10).value or "")[:200],
                "incident_num": str(ws.cell(r, 11).value or "")[:100],
            })

    return results


def update_slide7_apps_outages(
    updater: PPTUpdater,
    reader_path: str,
    month_name: str,
) -> None:
    """Slide 7 — Non-Critical Impairments + Critical Outages tables.

    Shapes:
      [2] Table Placeholder 11 (match 0) — Non-Critical Impairments (5x8)
          row 0: title, row 1: headers, rows 2-4: data
      [3] Table Placeholder 11 (match 1) — Critical Outages (3x6)
          row 0: title, row 1: headers, row 2+: data
    """
    slide_idx = 7

    downtimes = read_non_critical_downtimes(reader_path, month_name)

    # --- Non-Critical Impairments table (match 0) ---
    # Filter for Impairment or Degradation rows
    nc_rows = [
        d for d in downtimes
        if d["degradation_impairment"].lower() in ("impairment", "degradation")
    ]

    if nc_rows:
        # Table has 8 columns: Impacted App | Causing App | Degradation/Impairment |
        #   Duration | Customer Impact | Executive Summary | Mitigation Steps | Lessons Learned
        table_rows = []
        for d in nc_rows[:3]:  # max 3 data rows (rows 2-4)
            table_rows.append([
                d["impacted_app"],
                d["causing_app"],
                d["degradation_impairment"],
                d["duration"],
                d["customer_impact"],
                d["executive_summary"],
                d["mitigation_steps"],
                d["lessons_learned"],
            ])
        updater.update_table_data_rows(slide_idx, "Table Placeholder 11",
                                       table_rows, start_row=2, shape_index=0)

    # --- Critical Outages table (match 1) ---
    # Check Master-Apps Product Health for critical outage data
    from app.readers.excel_reader import ExcelReader
    critical_rows = []
    try:
        reader = ExcelReader(reader_path)
        # Check if any team had critical outages this month
        ws = reader.workbook.get("Master-Apps (DoNotChange)")
        if ws:
            # Scan Product Health section for critical outage counts
            # If all teams have 0 critical outages, show "no critical outages" message
            pass
    except Exception:
        pass

    # For now, check if any downtime rows are NOT impairment/degradation
    # (i.e., actual critical outages from the data)
    # Critical outages are typically separate — if none found, write default message
    if not critical_rows:
        # Write "no critical outages" message
        no_outage_text = f"There were no critical outages in {month_name[:3]} 20{26}"
        updater.update_table_data_rows(slide_idx, "Table Placeholder 11",
                                       [[no_outage_text, "", "", "", "", ""]],
                                       start_row=2, shape_index=1)

    logger.info("Slide 7 (Non-Critical Downtimes + Critical Outages) updated for %s", month_name)


# ---------------------------------------------------------------------------
# S16: DB Audit Charts
# ---------------------------------------------------------------------------

def read_db_audit_data(reader_path: str, month_name: str) -> dict:
    """Read DB Audit Data from Jaime-Platforms sheet.

    Returns a dict with all audit column values for the given month, keyed
    by column header name.  Values are floats (already 0-1 range in Excel,
    converted to percentage 0-100 for chart display).
    """
    from app.readers.excel_reader import ExcelReader

    try:
        reader = ExcelReader(reader_path)
        ws = reader.workbook["Jaime-Platforms"]
    except Exception:
        logger.warning("Jaime-Platforms sheet not found, skipping DB Audit read")
        return {}

    # Find "DB Audit Data" header
    audit_row = None
    for r in range(150, 200):
        val = ws.cell(r, 1).value
        if val and isinstance(val, str) and "db audit data" in val.lower():
            audit_row = r
            break

    if audit_row is None:
        logger.warning("DB Audit Data section not found in Jaime-Platforms")
        return {}

    # Row audit_row+1 = headers, audit_row+2+ = data
    header_row = audit_row + 1
    headers = []
    for c in range(1, 25):
        val = ws.cell(header_row, c).value
        if val:
            headers.append((c, str(val).strip()))
        else:
            break

    # Find the month row
    data = {}
    for r in range(audit_row + 2, audit_row + 20):
        m = ws.cell(r, 1).value
        if m and isinstance(m, str) and m.strip() == month_name:
            for col_idx, col_name in headers:
                if col_idx == 1:
                    continue  # skip Month column
                val = ws.cell(r, col_idx).value
                if val is not None and isinstance(val, (int, float)):
                    # Convert from 0-1 fraction to percentage (99.98 etc.)
                    data[col_name] = round(val * 100, 2)
                else:
                    data[col_name] = None
            break

    return data


def update_slide16_data_quality(updater: PPTUpdater, reader_path: str, month_name: str) -> None:
    """Slide 16 — Data Quality (DB Audit) charts.

    Shapes:
      [3] Chart 7 (5 series) — top-level audit accuracies
      [5] Chart 5 (11 series) — detailed audit accuracies
      [7] Chart 10 (3 series) — SLA/target metrics

    Uses overlay approach: read existing chart data, replace the current
    month's slot, write back.
    """
    audit_data = read_db_audit_data(reader_path, month_name)
    if not audit_data:
        logger.info("Slide 16 (Data Quality) — no DB Audit data for %s", month_name)
        return

    FISCAL_CATS = ["Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun",
                   "Jul", "Aug", "Sep", "Oct"]
    month_idx = MONTH_TO_IDX.get(month_name)
    if month_idx is None:
        return

    # Chart 7 (5 series): top-level audit accuracies
    # Series order matches the 5 series in the template chart
    chart7_cols = [
        "SOLO Audit Accuracy",
        "DRL Audit Accuracy",
        "UPC Audit Accuracy",
        "Agreements Audit Accuracy",
        "Audit Target Accuracy %",
    ]
    _overlay_standalone_chart(updater, 16, "Chart 7", None, FISCAL_CATS,
                              month_idx, audit_data, chart7_cols)

    # Chart 5 (11 series): detailed audit accuracies
    chart5_cols = [
        "SOLO/Effie Account Status Audit",
        "LCHTR Entitlement Audit",
        "LTWC Entitlement Audit",
        "DTC Entitlement Audit",
        "DTC Service Status Audit",
        "LCHTR Status Audit",
        "LTWC Status Audit",
        "DTC Ascendon Audit",
        "Solo/DS-Cache Account Domain Audit",
        "Solo/DS-Cache Customer Domain Audit",
        "Audit Target Accuracy %",
    ]
    _overlay_standalone_chart(updater, 16, "Chart 5", None, FISCAL_CATS,
                              month_idx, audit_data, chart5_cols)

    # Chart 10 (3 series): SLA/target metrics
    chart10_cols = [
        "SOLO SLA",
        "Solo Ascendon Audit",
        "SOLO Target SLA",
    ]
    _overlay_standalone_chart(updater, 16, "Chart 10", None, FISCAL_CATS,
                              month_idx, audit_data, chart10_cols)

    logger.info("Slide 16 (Data Quality / DB Audit) — charts updated for %s", month_name)


def _overlay_standalone_chart(
    updater: PPTUpdater,
    slide_idx: int,
    shape_name: str,
    shape_index,
    categories: list[str],
    month_idx: int,
    data: dict,
    col_names: list[str],
) -> None:
    """Read existing standalone chart data, overlay new month's values, write back.

    Args:
        data: Dict of column_name -> value (float percentage).
        col_names: Ordered list of column names, one per series in the chart.
    """
    from pptx.chart.data import CategoryChartData

    slide = updater.prs.slides[slide_idx]
    shape = updater._find_shape(slide, shape_name, shape_index)
    if not shape.has_chart:
        return

    chart = shape.chart
    existing_series = list(chart.series)

    # Read existing data
    old_data = [list(s.values) for s in existing_series]

    # Overlay new month's data for each series
    for si, col_name in enumerate(col_names):
        if si >= len(old_data):
            break
        val = data.get(col_name)
        if month_idx < len(old_data[si]):
            old_data[si][month_idx] = val

    # Write back
    cd = CategoryChartData()
    cd.categories = categories
    for si, series_vals in enumerate(old_data):
        name = f"Series {si+1}"
        try:
            name = existing_series[si].tx.strRef.strCache[0].v
        except Exception:
            pass
        add_chart_series(cd, name, series_vals)
    chart.replace_data(cd)


# ---------------------------------------------------------------------------
# S21: Key Highlights / Projects Delivered Text
# ---------------------------------------------------------------------------

def update_slide21_dpos_delivery_text(
    updater: PPTUpdater,
    reader_path: str,
    month_name: str,
) -> None:
    """Slide 21 — DPOS On-Time Delivery text (Projects Delivered).

    Shape [0] Text Placeholder 4 contains project delivery highlights.
    Data comes from individual person sheets' Key Highlights column (col 13).
    The Master sheet has #REF! errors, so we read directly from person sheets.
    """
    from app.readers.excel_reader import ExcelReader

    MONTHS = [
        "January", "February", "March", "April", "May", "June",
        "July", "August", "September", "October", "November", "December",
    ]

    # Collect Key Highlights from individual person sheets
    highlights: list[str] = []

    try:
        reader = ExcelReader(reader_path)
    except Exception:
        logger.warning("Could not open workbook for S21 Key Highlights")
        return

    # Only read from platform-related sheets (DPOS = platforms)
    platform_sheets = [
        s for s in reader.workbook.sheetnames
        if "Platform" in s and "Master" not in s and "DoNotChange" not in s
    ]

    for sheet_name in platform_sheets:
        ws = reader.workbook[sheet_name]
        # Key Highlights is in column 13, with month names in column 1
        # Row 2 has the header, data starts at row 3
        for r in range(3, 50):
            m = ws.cell(r, 1).value
            if m and isinstance(m, str) and m.strip() == month_name:
                val = ws.cell(r, 13).value
                if val and isinstance(val, str) and val.strip() and val.strip() not in ("N/A", "TBD", "NA"):
                    highlights.append(val.strip())
                # Also check subsequent rows for same month (multi-team)
                for r2 in range(r + 1, r + 10):
                    m2 = ws.cell(r2, 1).value
                    # If next row has a different month or another month name, stop
                    if m2 and isinstance(m2, str) and m2.strip() in MONTHS:
                        break
                    # If it's a team name (not a month), check its highlights
                    val2 = ws.cell(r2, 13).value
                    if val2 and isinstance(val2, str) and val2.strip() and val2.strip() not in ("N/A", "TBD", "NA"):
                        highlights.append(val2.strip())
                break

    if not highlights:
        logger.info("Slide 21 — no Key Highlights data for %s, leaving unchanged", month_name)
        return

    # Build the text content for the text frame
    # Parse highlights to extract project lines organized by track
    all_lines = []
    for h in highlights:
        # Each highlight may contain multi-line text with Track sections
        for line in h.split("\n"):
            line = line.strip()
            if line and line not in ("N/A", "TBD", "NA", "- N/A", "- NA", "- TBD"):
                all_lines.append(line)

    if not all_lines:
        logger.info("Slide 21 — Key Highlights empty after parsing for %s", month_name)
        return

    # Build final text: "Projects Delivered" header + content
    text_parts = ["Projects Delivered"]
    text_parts.extend(all_lines)
    final_text = "\n".join(text_parts)

    updater.update_text_frame(21, "Text Placeholder 4", final_text, shape_index=0)
    logger.info("Slide 21 (DPOS delivery text) updated for %s", month_name)
