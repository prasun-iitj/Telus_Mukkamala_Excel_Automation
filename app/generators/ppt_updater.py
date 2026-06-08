"""PPT Updater – clone template and fill with data using python-pptx."""

from __future__ import annotations

import logging
import shutil
from copy import deepcopy
from typing import Optional, Union

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.util import Pt

from app.exceptions import ShapeNotFoundError
from app.generators.chart_format import add_chart_series
from app.models import ChartDataSpec, TableDataSpec

logger = logging.getLogger(__name__)


class PPTUpdater:
    """Clone a PowerPoint template and update shapes in-place.

    The template file is never modified – a copy is made at *output_path*
    and all mutations happen on the copy.
    """

    def __init__(self, template_path: str, output_path: str) -> None:
        """Copy *template_path* to *output_path*, then open the copy.

        Args:
            template_path: Path to the original .pptx template.
            output_path: Destination path for the generated .pptx.
        """
        self.template_path = template_path
        self.output_path = output_path
        shutil.copy2(template_path, output_path)
        self.prs = Presentation(output_path)

    # ------------------------------------------------------------------
    # Shape lookup
    # ------------------------------------------------------------------

    def _find_shape(
        self,
        slide,
        shape_name: str,
        shape_index: Optional[int] = None,
    ):
        """Locate a shape on *slide* by name, with match-index disambiguation.

        Strategy:
        1. Collect all shapes whose ``name`` matches *shape_name*.
        2. If exactly one match → return it.
        3. If multiple matches and *shape_index* is given → use it as
           an index into the *matches* list (Nth shape with that name).
        4. If multiple matches but no index → return the first one.
        5. If no name match → raise :class:`ShapeNotFoundError`.
        """
        matches = [s for s in slide.shapes if s.name == shape_name]

        if len(matches) == 1:
            return matches[0]

        if len(matches) > 1:
            if shape_index is not None and 0 <= shape_index < len(matches):
                return matches[shape_index]
            return matches[0]

        # Determine slide index for the error message
        slide_index = -1
        for idx, s in enumerate(self.prs.slides):
            if s == slide:
                slide_index = idx
                break
        raise ShapeNotFoundError(slide_index, shape_name)

    @staticmethod
    def _element_shape_name(element) -> Optional[str]:
        """Return the drawing name for a top-level slide shape element."""
        cnv = element.find(f".//{qn('p:cNvPr')}")
        if cnv is not None:
            return cnv.get("name")
        return None

    def reset_shapes_from_template(
        self,
        slide_index: int,
        shape_names: list[str],
        *,
        remove_name_substrings: Optional[list[str]] = None,
    ) -> None:
        """Replace named shapes on a slide with pristine copies from the template.

        Preserves z-order by inserting clones where the first removed shape was.
        Also removes any shapes whose name contains a substring in
        *remove_name_substrings* (e.g. MoM arrows).
        """
        tpl = Presentation(self.template_path)
        tpl_slide = tpl.slides[slide_index]
        out_slide = self.prs.slides[slide_index]
        sp_tree = out_slide.shapes._spTree

        names_to_reset = set(shape_names)
        substrings = [s.lower() for s in (remove_name_substrings or [])]

        def _should_remove(element) -> bool:
            name = self._element_shape_name(element)
            if not name:
                return False
            if name in names_to_reset:
                return True
            return any(sub in name.lower() for sub in substrings)

        children = list(sp_tree)
        insert_idx = next((i for i, el in enumerate(children) if _should_remove(el)), None)

        to_copy = [
            deepcopy(child)
            for child in tpl_slide.shapes._spTree
            if self._element_shape_name(child) in names_to_reset
        ]

        for i in range(len(children) - 1, -1, -1):
            if _should_remove(children[i]):
                sp_tree.remove(children[i])

        if not to_copy:
            return

        if insert_idx is None:
            insert_idx = len(list(sp_tree))

        for offset, element in enumerate(to_copy):
            sp_tree.insert(insert_idx + offset, element)

    def _purge_arrow_shapes(self, slide_index: int) -> None:
        """Remove MoM arrow shapes only (they overlap table rows when data is N/A)."""
        sp_tree = self.prs.slides[slide_index].shapes._spTree
        to_remove = []
        for child in sp_tree:
            name = self._element_shape_name(child) or ""
            prst_geoms = child.findall(f".//{qn('a:prstGeom')}")
            is_arrow = "arrow" in name.lower() or any(
                "arrow" in (geom.get("prst") or "").lower() for geom in prst_geoms
            )
            if is_arrow:
                to_remove.append(child)
        for child in to_remove:
            sp_tree.remove(child)

    # ------------------------------------------------------------------
    # Chart update
    # ------------------------------------------------------------------

    def update_chart(
        self,
        slide_index: int,
        shape_name: str,
        chart_data: ChartDataSpec,
        shape_index: Optional[int] = None,
    ) -> None:
        """Replace chart data while preserving all chart formatting.

        None values in a series are passed through to python-pptx which
        renders them as gaps (not zero).

        Args:
            slide_index: 0-based slide position.
            shape_name: Name of the chart shape.
            chart_data: A :class:`ChartDataSpec` with categories and series.
            shape_index: Optional positional disambiguator.
        """
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)

        if not shape.has_chart:
            logger.warning(
                "Shape '%s' on slide %d does not contain a chart — skipping chart update",
                shape_name,
                slide_index,
            )
            return

        cd = CategoryChartData()
        cd.categories = chart_data.categories

        for series in chart_data.series:
            # None values are kept as-is; python-pptx treats them as gaps
            add_chart_series(cd, series.name, series.values)

        shape.chart.replace_data(cd)

    def update_grouped_chart(
        self,
        slide_index: int,
        group_index: int,
        chart_index: int,
        chart_data: ChartDataSpec,
    ) -> None:
        """Update a chart that's inside a group shape.

        Args:
            slide_index: 0-based slide position.
            group_index: Index of the group shape among all shapes on the slide.
            chart_index: Index of the chart shape within the group.
            chart_data: A :class:`ChartDataSpec` with categories and series.
        """
        slide = self.prs.slides[slide_index]
        shapes_list = list(slide.shapes)
        if group_index >= len(shapes_list):
            logger.warning("Group index %d out of range on slide %d", group_index, slide_index)
            return

        group = shapes_list[group_index]
        if group.shape_type != 6:  # MSO_SHAPE_TYPE.GROUP
            logger.warning("Shape [%d] on slide %d is not a group", group_index, slide_index)
            return

        group_shapes = list(group.shapes)
        if chart_index >= len(group_shapes):
            logger.warning("Chart index %d out of range in group on slide %d", chart_index, slide_index)
            return

        child = group_shapes[chart_index]
        if not hasattr(child, "has_chart") or not child.has_chart:
            logger.warning("Shape [%d] in group [%d] on slide %d has no chart", chart_index, group_index, slide_index)
            return

        cd = CategoryChartData()
        cd.categories = chart_data.categories
        for series in chart_data.series:
            add_chart_series(cd, series.name, series.values)
        child.chart.replace_data(cd)

    # ------------------------------------------------------------------
    # Table update
    # ------------------------------------------------------------------

    def update_table(
        self,
        slide_index: int,
        shape_name: str,
        table_data: TableDataSpec,
        shape_index: Optional[int] = None,
    ) -> None:
        """Update table cell text while preserving cell formatting.

        If the first paragraph of a cell already has a run, only the run
        text is replaced so that font/color/size are kept.  Otherwise
        ``cell.text`` is set directly.

        Rows that exceed the existing table size are silently skipped
        with a warning.

        Args:
            slide_index: 0-based slide position.
            shape_name: Name of the table shape.
            table_data: A :class:`TableDataSpec` with headers and rows.
            shape_index: Optional positional disambiguator.
        """
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)

        if not shape.has_table:
            logger.warning(
                "Shape '%s' on slide %d does not contain a table — skipping table update",
                shape_name,
                slide_index,
            )
            return

        table = shape.table

        # Determine header offset – if headers are provided, first row is headers
        header_offset = 1 if table_data.headers else 0

        # Update header row if provided
        if table_data.headers:
            for col_idx, header in enumerate(table_data.headers):
                if col_idx < len(table.columns):
                    cell = table.cell(0, col_idx)
                    self._set_cell_text(cell, str(header) if header is not None else "")

        available_rows = len(table.rows) - header_offset
        if len(table_data.rows) > available_rows:
            logger.warning(
                "Table on slide %d, shape '%s': data has %d rows but table "
                "only has %d data rows. Extra rows will be skipped.",
                slide_index,
                shape_name,
                len(table_data.rows),
                available_rows,
            )

        for row_idx, row_data in enumerate(table_data.rows):
            target_row = row_idx + header_offset
            if target_row >= len(table.rows):
                break
            for col_idx, value in enumerate(row_data):
                if col_idx >= len(table.columns):
                    break
                cell = table.cell(target_row, col_idx)
                text_val = str(value) if value is not None else ""
                self._set_cell_text(cell, text_val)

    def update_table_data_rows(
        self,
        slide_index: int,
        shape_name: str,
        data_rows: list[list],
        start_row: int = 3,
        shape_index: Optional[int] = None,
        data_start_col: int = 0,
    ) -> None:
        """Update only the data rows of a table, leaving title/summary/headers untouched.

        Args:
            slide_index: 0-based slide position.
            shape_name: Name of the table shape.
            data_rows: List of rows to write. Each row is a list of cell values.
            start_row: 0-based row index to start writing (default 3, skipping title/summary/headers).
            shape_index: Optional Nth-match disambiguator.
            data_start_col: Column index where the first value in each row is written.
        """
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)

        if not shape.has_table:
            logger.warning(
                "Shape '%s' on slide %d is not a table — skipping",
                shape_name, slide_index,
            )
            return

        table = shape.table
        for row_idx, row_data in enumerate(data_rows):
            target_row = start_row + row_idx
            if target_row >= len(table.rows):
                logger.warning(
                    "Table on slide %d '%s': row %d out of range (%d rows total)",
                    slide_index, shape_name, target_row, len(table.rows),
                )
                break
            for col_idx, value in enumerate(row_data):
                target_col = data_start_col + col_idx
                if target_col >= len(table.columns):
                    break
                cell = table.cell(target_row, target_col)
                text_val = str(value) if value is not None else ""
                self._set_cell_text(cell, text_val)

    def update_table_by_team_match(
        self,
        slide_index: int,
        shape_name: str,
        month_data,
        row_builder,
        start_row: int = 3,
        shape_index: Optional[int] = None,
        data_start_col: int = 1,
    ) -> int:
        """Update table rows by matching PPT team labels to Excel data.

        Preserves existing team-name cells (column 0) and only writes data
        columns. Returns the number of rows updated.

        Args:
            month_data: MonthData with team metrics from Excel.
            row_builder: Callable(TeamMetrics) -> list of cell values (no team name).
            data_start_col: Column index where data values begin (default 1).
        """
        from app.config.team_matching import find_team_for_ppt_row

        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)

        if not shape.has_table:
            logger.warning(
                "Shape '%s' on slide %d is not a table — skipping",
                shape_name, slide_index,
            )
            return 0

        table = shape.table
        updated = 0

        for row_idx in range(start_row, len(table.rows)):
            ppt_team = table.cell(row_idx, 0).text.strip()
            if not ppt_team:
                continue

            team_metrics = find_team_for_ppt_row(ppt_team, month_data)
            if team_metrics is None:
                continue

            row_data = row_builder(team_metrics)
            for col_offset, value in enumerate(row_data):
                col_idx = data_start_col + col_offset
                if col_idx >= len(table.columns):
                    break
                cell = table.cell(row_idx, col_idx)
                text_val = str(value) if value is not None else ""
                self._set_cell_text(cell, text_val)
            updated += 1

        return updated

    def update_table_by_dpos_codename(
        self,
        slide_index: int,
        shape_name: str,
        month_data,
        dpos_map,
        row_builder,
        start_row: int = 3,
        shape_index: Optional[int] = None,
        data_start_col: int = 1,
    ) -> int:
        """Update DPOS detail table rows using codename → platform team mapping.

        Reads the codename prefix from column 0 (e.g. 'Byron' from
        'Byron\\n(Account, Identity, Agreements)'), resolves the mapped
        platform team, and writes data columns without changing team labels.
        """
        from app.config.dpos_team_map_parser import extract_dpos_codename
        from app.config.team_matching import find_team_by_platform_name

        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)

        if not shape.has_table:
            logger.warning(
                "Shape '%s' on slide %d is not a table — skipping",
                shape_name, slide_index,
            )
            return 0

        table = shape.table
        updated = 0

        for row_idx in range(start_row, len(table.rows)):
            ppt_team_cell = table.cell(row_idx, 0).text.strip()
            if not ppt_team_cell:
                continue

            codename = extract_dpos_codename(ppt_team_cell)
            mapping = dpos_map.teams.get(codename)
            if mapping is None:
                logger.debug(
                    "No DPOS mapping for codename '%s' on slide %d",
                    codename, slide_index,
                )
                continue

            team_metrics = find_team_by_platform_name(
                mapping.platform_team, month_data,
            )
            if team_metrics is None:
                logger.debug(
                    "Platform team '%s' not found for DPOS codename '%s'",
                    mapping.platform_team, codename,
                )
                continue

            row_data = row_builder(team_metrics)
            for col_offset, value in enumerate(row_data):
                col_idx = data_start_col + col_offset
                if col_idx >= len(table.columns):
                    break
                cell = table.cell(row_idx, col_idx)
                text_val = str(value) if value is not None else ""
                self._set_cell_text(cell, text_val)
            updated += 1

        return updated

    @staticmethod
    def _set_cell_text(cell, text: str) -> None:
        """Set cell text preserving existing run formatting."""
        cell.text_frame.word_wrap = True
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = text
            for ri in range(len(cell.text_frame.paragraphs[0].runs) - 1, 0, -1):
                cell.text_frame.paragraphs[0].runs[ri].text = ""
        else:
            cell.text = text

    # ------------------------------------------------------------------
    # Text frame update
    # ------------------------------------------------------------------

    def update_text_frame(
        self,
        slide_index: int,
        shape_name: str,
        text_content: Union[str, list[str]],
        shape_index: Optional[int] = None,
    ) -> None:
        """Update a text frame preserving paragraph/run formatting.

        If *text_content* is a string it is split on newlines.  Each line
        maps to an existing paragraph; extra lines beyond the paragraph
        count are logged as a warning and skipped.

        Args:
            slide_index: 0-based slide position.
            shape_name: Name of the text-frame shape.
            text_content: Text to set – a string (may contain ``\\n``) or
                a list of strings (one per paragraph).
            shape_index: Optional positional disambiguator.
        """
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)
        text_frame = shape.text_frame

        if text_content is None or text_content == "":
            # Clear all paragraphs
            for para in text_frame.paragraphs:
                if para.runs:
                    for r in para.runs:
                        r.text = ""
                else:
                    para.text = ""
            return

        lines: list[str] = (
            text_content.split("\n")
            if isinstance(text_content, str)
            else list(text_content)
        )

        if len(lines) > len(text_frame.paragraphs):
            logger.warning(
                "Text frame on slide %d, shape '%s': content has %d lines "
                "but text frame only has %d paragraphs. Extra lines skipped.",
                slide_index,
                shape_name,
                len(lines),
                len(text_frame.paragraphs),
            )

        for para_idx, content_line in enumerate(lines):
            if para_idx >= len(text_frame.paragraphs):
                break
            para = text_frame.paragraphs[para_idx]
            if para.runs:
                para.runs[0].text = content_line
                # Clear any additional runs to prevent text concatenation
                for ri in range(len(para.runs) - 1, 0, -1):
                    para.runs[ri].text = ""
            else:
                para.text = content_line

    def hide_shape(self, slide_index: int, shape_name: str) -> None:
        """Hide a shape on a slide (e.g. MoM arrow icons when data is N/A)."""
        from pptx.oxml.ns import qn

        slide = self.prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == shape_name:
                c_nv_pr = shape.element.xpath(".//p:cNvPr")
                if c_nv_pr:
                    c_nv_pr[0].set("hidden", "1")
                return

    def hide_shapes_matching(self, slide_index: int, name_substring: str) -> None:
        """Hide all shapes whose name contains *name_substring*."""
        slide = self.prs.slides[slide_index]
        for shape in slide.shapes:
            if name_substring.lower() in shape.name.lower():
                c_nv_pr = shape.element.xpath(".//p:cNvPr")
                if c_nv_pr:
                    c_nv_pr[0].set("hidden", "1")

    def remove_shapes_matching(self, slide_index: int, name_substring: str) -> None:
        """Remove shapes whose name contains *name_substring* (e.g. MoM arrows)."""
        slide = self.prs.slides[slide_index]
        to_remove = [
            shape for shape in slide.shapes
            if name_substring.lower() in shape.name.lower()
        ]
        for shape in to_remove:
            shape.element.getparent().remove(shape.element)

    def update_table_cell(
        self,
        slide_index: int,
        shape_name: str,
        row: int,
        col: int,
        value: str,
        shape_index: Optional[int] = None,
    ) -> None:
        """Update a single table cell without rewriting the whole row."""
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)
        if not shape.has_table:
            return
        table = shape.table
        if row < len(table.rows) and col < len(table.columns):
            self._set_cell_text(table.cell(row, col), str(value) if value is not None else "")

    def set_table_row_heights(
        self,
        slide_index: int,
        shape_name: str,
        row_heights: dict[int, int],
        shape_index: Optional[int] = None,
    ) -> None:
        """Set explicit row heights (EMU) for selected rows."""
        slide = self.prs.slides[slide_index]
        shape = self._find_shape(slide, shape_name, shape_index)
        if not shape.has_table:
            return
        table = shape.table
        for row_idx, height in row_heights.items():
            if row_idx < len(table.rows):
                table.rows[row_idx].height = height

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------

    def save(self) -> str:
        """Save the presentation and return the output file path."""
        self.prs.save(self.output_path)
        return self.output_path
