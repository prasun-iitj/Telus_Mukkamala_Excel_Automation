"""Unit tests for PPTUpdater class."""

import os
import shutil
import tempfile

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from app.exceptions import ShapeNotFoundError
from app.generators.ppt_updater import PPTUpdater
from app.models import ChartDataSpec, SeriesSpec, TableDataSpec


# ---------------------------------------------------------------------------
# Helpers – create minimal .pptx templates for testing
# ---------------------------------------------------------------------------

def _create_template_with_text(path: str) -> None:
    """Create a minimal .pptx with a single text shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.name = "TestTextBox"
    tf = txBox.text_frame
    tf.text = "Line 1"
    p = tf.add_paragraph()
    p.text = "Line 2"
    p2 = tf.add_paragraph()
    p2.text = "Line 3"
    prs.save(path)


def _create_template_with_table(path: str, rows: int = 4, cols: int = 3) -> None:
    """Create a minimal .pptx with a table shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(3))
    table_shape.name = "TestTable"
    table = table_shape.table
    # Fill with placeholder data
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = f"R{r}C{c}"
    prs.save(path)


def _create_template_with_chart(path: str) -> None:
    """Create a minimal .pptx with a chart shape."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Jan", "Feb", "Mar"]
    chart_data.add_series("Series A", (1, 2, 3))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(6), Inches(4),
        chart_data,
    )
    chart_shape.name = "TestChart"
    prs.save(path)


def _create_template_with_duplicate_names(path: str) -> None:
    """Create a .pptx with multiple shapes sharing the same name."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(3):
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1 + i), Inches(4), Inches(0.5)
        )
        txBox.name = "DuplicateName"
        txBox.text_frame.text = f"Box {i}"
    prs.save(path)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def tmp_dir():
    d = tempfile.mkdtemp()
    yield d
    shutil.rmtree(d, ignore_errors=True)


# ---------------------------------------------------------------------------
# Tests: __init__ / template cloning
# ---------------------------------------------------------------------------

class TestInit:
    def test_template_is_copied(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        assert os.path.exists(out)
        # Template should be untouched
        assert os.path.getsize(tpl) == os.path.getsize(out)

    def test_presentation_is_loaded(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        assert updater.prs is not None
        assert len(updater.prs.slides) == 1


# ---------------------------------------------------------------------------
# Tests: _find_shape
# ---------------------------------------------------------------------------

class TestFindShape:
    def test_find_by_name(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)
        updater = PPTUpdater(tpl, out)
        slide = updater.prs.slides[0]
        shape = updater._find_shape(slide, "TestTextBox")
        assert shape.name == "TestTextBox"

    def test_find_by_index_fallback(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)
        updater = PPTUpdater(tpl, out)
        slide = updater.prs.slides[0]
        shape = updater._find_shape(slide, "NonExistent", shape_index=0)
        assert shape is not None

    def test_raises_when_not_found(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)
        updater = PPTUpdater(tpl, out)
        slide = updater.prs.slides[0]
        with pytest.raises(ShapeNotFoundError):
            updater._find_shape(slide, "NonExistent")

    def test_duplicate_names_with_index(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_duplicate_names(tpl)
        updater = PPTUpdater(tpl, out)
        slide = updater.prs.slides[0]
        shape = updater._find_shape(slide, "DuplicateName", shape_index=1)
        assert shape is not None

    def test_duplicate_names_without_index(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_duplicate_names(tpl)
        updater = PPTUpdater(tpl, out)
        slide = updater.prs.slides[0]
        # Should return the first match
        shape = updater._find_shape(slide, "DuplicateName")
        assert shape.name == "DuplicateName"


# ---------------------------------------------------------------------------
# Tests: update_chart
# ---------------------------------------------------------------------------

class TestUpdateChart:
    def test_chart_data_replaced(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_chart(tpl)

        updater = PPTUpdater(tpl, out)
        spec = ChartDataSpec(
            categories=["Apr", "May", "Jun"],
            series=[SeriesSpec(name="New Series", values=[10, 20, 30])],
        )
        updater.update_chart(0, "TestChart", spec)
        updater.save()

        # Reload and verify
        prs = Presentation(out)
        chart = prs.slides[0].shapes[0].chart
        assert len(list(chart.series)) == 1

    def test_chart_with_none_values(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_chart(tpl)

        updater = PPTUpdater(tpl, out)
        spec = ChartDataSpec(
            categories=["Jan", "Feb", "Mar"],
            series=[SeriesSpec(name="S1", values=[10, None, 30])],
        )
        # Should not raise
        updater.update_chart(0, "TestChart", spec)
        updater.save()


# ---------------------------------------------------------------------------
# Tests: update_table
# ---------------------------------------------------------------------------

class TestUpdateTable:
    def test_table_cells_updated(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_table(tpl, rows=3, cols=2)

        updater = PPTUpdater(tpl, out)
        spec = TableDataSpec(
            headers=["H1", "H2"],
            rows=[["A", "B"], ["C", "D"]],
        )
        updater.update_table(0, "TestTable", spec)
        updater.save()

        prs = Presentation(out)
        table = prs.slides[0].shapes[0].table
        assert table.cell(0, 0).text == "H1"
        assert table.cell(0, 1).text == "H2"
        assert table.cell(1, 0).text == "A"
        assert table.cell(1, 1).text == "B"
        assert table.cell(2, 0).text == "C"
        assert table.cell(2, 1).text == "D"

    def test_table_with_none_values(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_table(tpl, rows=2, cols=2)

        updater = PPTUpdater(tpl, out)
        spec = TableDataSpec(headers=[], rows=[[None, "val"]])
        updater.update_table(0, "TestTable", spec)
        updater.save()

        prs = Presentation(out)
        table = prs.slides[0].shapes[0].table
        assert table.cell(0, 0).text == ""
        assert table.cell(0, 1).text == "val"

    def test_table_extra_rows_warning(self, tmp_dir):
        """More data rows than table rows should log warning, not crash."""
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_table(tpl, rows=2, cols=2)

        updater = PPTUpdater(tpl, out)
        spec = TableDataSpec(
            headers=["H1", "H2"],
            rows=[["A", "B"], ["C", "D"], ["E", "F"]],  # 3 data rows, only 1 fits
        )
        # Should not raise
        updater.update_table(0, "TestTable", spec)
        updater.save()


# ---------------------------------------------------------------------------
# Tests: update_text_frame
# ---------------------------------------------------------------------------

class TestUpdateTextFrame:
    def test_text_updated(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        updater.update_text_frame(0, "TestTextBox", "New Line 1\nNew Line 2\nNew Line 3")
        updater.save()

        prs = Presentation(out)
        tf = prs.slides[0].shapes[0].text_frame
        assert tf.paragraphs[0].text == "New Line 1"
        assert tf.paragraphs[1].text == "New Line 2"
        assert tf.paragraphs[2].text == "New Line 3"

    def test_text_as_list(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        updater.update_text_frame(0, "TestTextBox", ["A", "B", "C"])
        updater.save()

        prs = Presentation(out)
        tf = prs.slides[0].shapes[0].text_frame
        assert tf.paragraphs[0].text == "A"
        assert tf.paragraphs[1].text == "B"
        assert tf.paragraphs[2].text == "C"

    def test_empty_text_clears(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        updater.update_text_frame(0, "TestTextBox", "")
        updater.save()

    def test_extra_lines_warning(self, tmp_dir):
        """More lines than paragraphs should not crash."""
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)  # 3 paragraphs

        updater = PPTUpdater(tpl, out)
        updater.update_text_frame(0, "TestTextBox", "1\n2\n3\n4\n5")
        updater.save()


# ---------------------------------------------------------------------------
# Tests: save
# ---------------------------------------------------------------------------

class TestSave:
    def test_save_returns_path(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        updater = PPTUpdater(tpl, out)
        result = updater.save()
        assert result == out
        assert os.path.exists(out)

    def test_template_unchanged_after_save(self, tmp_dir):
        tpl = os.path.join(tmp_dir, "template.pptx")
        out = os.path.join(tmp_dir, "output.pptx")
        _create_template_with_text(tpl)

        # Read original template bytes
        with open(tpl, "rb") as f:
            original_bytes = f.read()

        updater = PPTUpdater(tpl, out)
        updater.update_text_frame(0, "TestTextBox", "Modified text")
        updater.save()

        # Template should be byte-identical
        with open(tpl, "rb") as f:
            after_bytes = f.read()
        assert original_bytes == after_bytes
